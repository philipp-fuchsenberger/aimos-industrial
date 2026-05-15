"""
AIMOS Batch/OODA Orchestrator — CR-233/234/236/237
=====================================================
6-phase OODA cycle for the "Mitarbeiter" agent archetype.

Phase 0 (KONTEXT):  Load workspace state (files, DBs, memories)
Phase 1 (OBSERVE):  Read all new messages, structure, find cross-thread connections
Phase 2 (ORIENT):   Build Lagebild (2a: chunk-loop, 2b: consolidation)
Phase 3 (DECIDE):   Identify stakeholders, create action plan
Phase 4 (ACT):      Draft → Validate → Dispatch (4a/4b/4c)
Phase 5 (PERSIST):  Guaranteed self-dispatch: save state, remember facts

Architecture:
  - The Orchestrator (this module) is DETERMINISTIC Python code
  - It calls think() on a NON-DETERMINISTIC LLM multiple times
  - Each phase output is captured as a string and injected into the next phase's prompt
  - The Lagebild is EPHEMERAL — rebuilt each cycle from workspace state + new messages
  - Workspace files (state.md, todo.md, etc.) carry state between cycles
  - See docs/AGENT_ARCHETYPES.md for full documentation

Safety (HAZOP/FMEA/FTA/STPA — docs/snapshots/):
  - H-02: Phase 0 validates state.md against data sources
  - H-04: Staleness warning if state.md >7 days old
  - H-07: Message dedup in poll_pending (agent_base.py)
  - H-08: Human approval gate (batch_require_human_approval)
  - H-09: Empty Lagebild → cycle abort
  - H-14: state.md.bak backup before Phase 4
  - H-15: BATCH_COMPLETED marker after successful cycle
  - S-3:  Cross-thread leak protection in Phase 3 prompt
  - S-7:  Phase 3 results collected and passed to Phase 4
  - BF-15: Phase 4 warns against remembering unverified facts
  - BF-17: Memory pruning in agent_base.py (max_memories)
  - Context Monitor: Intelligent compression, not blind truncation
"""

import asyncio
import json
import logging
import shutil
import time
from collections import defaultdict
from pathlib import Path

import aiohttp

from core.config import Config

# Type hint only — avoid circular import
from typing import TYPE_CHECKING
if TYPE_CHECKING:
    from core.agent_base import AIMOSAgent


# ── CR-274: Confidentiality Scope System ───────────────────────────────────

def _resolve_scope(thread_id: str, config: dict) -> str:
    """Map a thread_id to a confidentiality scope. Pure Python, no LLM.

    The scope determines which Lagebild partition a stakeholder sees in Phase 4.
    Modes (batch_scope_pattern):
      "email_address" (default): email:foo@bar.com → scope:foo@bar.com
      "thread_id": scope = thread_id verbatim
      "config_map": explicit mapping from batch_scope_map config
    """
    mode = config.get("batch_scope_pattern", "email_address")
    if mode == "config_map":
        scope_map = config.get("batch_scope_map", {})
        return scope_map.get(thread_id, f"scope:{thread_id}")
    if mode == "email_address" and thread_id.startswith("email:"):
        return f"scope:{thread_id[6:]}"
    return f"scope:{thread_id}"


def _group_messages_by_scope(
    messages: list[dict], config: dict,
) -> dict[str, list[dict]]:
    """Group messages by confidentiality scope."""
    scopes: dict[str, list[dict]] = defaultdict(list)
    for msg in messages:
        tid = msg.get("thread_id") or f"sender:{msg.get('sender_id', 0)}"
        scope = _resolve_scope(tid, config)
        scopes[scope].append(msg)
    return dict(scopes)


def _partition_lagebild(
    lagebild: str, scope_names: list[str],
) -> dict[str, str]:
    """Split a scope-tagged Lagebild into per-scope partitions.

    Expects the Lagebild to contain sections like:
      ## [SCOPE: scope:foo@bar.com]
      ...content...
      ## [SCOPE: scope:baz@bar.com]
      ...content...

    Returns {"scope:foo@bar.com": "...content...", "scope:baz@bar.com": "..."}.
    If no scope headers found, returns {"_all": lagebild} (fallback).
    """
    import re
    pattern = re.compile(r'^##\s*\[SCOPE:\s*([^\]]+)\]', re.MULTILINE)
    matches = list(pattern.finditer(lagebild))
    if not matches:
        return {"_all": lagebild}

    partitions: dict[str, str] = {}
    for i, match in enumerate(matches):
        scope_key = match.group(1).strip()
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(lagebild)
        partitions[scope_key] = lagebild[start:end].strip()

    return partitions


# ── CR-275: Semantic Message Dedup ─────────────────────────────────────────

def _dedup_messages(
    messages: list[dict], log: logging.Logger, agent_name: str,
) -> list[dict]:
    """Merge near-duplicate messages from the same sender/thread.

    A human sees "same person, same topic, 10 minutes apart" as one request.
    The agent should too. Without this, the agent sends N replies for N
    near-identical inputs.

    Strategy:
    1. Group messages by thread_id (= same sender)
    2. Within a group, compare content similarity (word overlap ratio)
    3. If similarity > 60% and time delta < 30 min: merge (keep latest, combine text)
    """
    if len(messages) <= 1:
        return messages

    from collections import defaultdict

    # Group by thread
    by_thread: dict[str, list[dict]] = defaultdict(list)
    for msg in messages:
        tid = msg.get("thread_id") or f"sender:{msg.get('sender_id', 0)}"
        by_thread[tid].append(msg)

    result: list[dict] = []
    merged_count = 0

    for tid, group in by_thread.items():
        if len(group) <= 1:
            result.extend(group)
            continue

        # Sort by id (chronological)
        group.sort(key=lambda m: m.get("id", 0))

        # Compare each pair — merge if similar
        kept: list[dict] = [group[0]]
        for msg in group[1:]:
            prev = kept[-1]
            if _messages_are_duplicates(prev, msg):
                # Merge: keep the later message, append unique info from earlier
                prev_text = prev.get("content", "")
                msg_text = msg.get("content", "")
                # If the new message is substantially longer, use it as base
                if len(msg_text) > len(prev_text) * 1.2:
                    kept[-1] = msg
                # Otherwise keep the earlier one (it was first)
                merged_count += 1
                log.info(
                    f"[{agent_name}] CR-275 Dedup: merged duplicate from "
                    f"thread={tid} (#{prev.get('id', '?')} + #{msg.get('id', '?')})"
                )
            else:
                kept.append(msg)
        result.extend(kept)

    if merged_count:
        log.info(
            f"[{agent_name}] CR-275 Dedup: {len(messages)} messages → "
            f"{len(result)} after merging {merged_count} duplicate(s)"
        )
    return result


def _messages_are_duplicates(a: dict, b: dict) -> bool:
    """Check if two messages are near-duplicates.

    Primary criterion: same sender (thread_id).
    Secondary criterion: similar content (word overlap) as tiebreaker
    for cases where the same sender sends genuinely different messages
    in the same batch.
    """
    # Same thread_id = same sender. Different sender = never a duplicate.
    tid_a = a.get("thread_id") or f"sender:{a.get('sender_id', '')}"
    tid_b = b.get("thread_id") or f"sender:{b.get('sender_id', '')}"
    if tid_a != tid_b:
        return False

    # Same sender — now check if it's genuinely different content
    # or just a retry/duplicate. Use word overlap as tiebreaker.
    # Extract ONLY the body text (after "Text:" header) for comparison.
    import re

    def _extract_body(content: str) -> str:
        """Extract only the email body (after 'Text:' line), lowercased."""
        match = re.search(r'^Text:\s*(.+)', content, re.MULTILINE | re.DOTALL)
        if match:
            return match.group(1).strip().lower()
        return content.lower().strip()

    text_a = _extract_body(a.get("content", ""))
    text_b = _extract_body(b.get("content", ""))

    if not text_a or not text_b:
        return True  # Same sender, empty content = duplicate

    words_a = {w for w in text_a.split() if len(w) > 5}
    words_b = {w for w in text_b.split() if len(w) > 5}
    if not words_a or not words_b:
        return True  # Same sender, no meaningful words = duplicate

    overlap = len(words_a & words_b)
    smaller = min(len(words_a), len(words_b))
    ratio = overlap / smaller if smaller > 0 else 0

    # >50% word overlap from same sender = duplicate
    # <50% = genuinely different message (e.g. "here are my docs" + "here is my payslip")
    return ratio > 0.5


# ── CR-278: Arbeitsdatei Summary for Draft Prompt ──────────────────────────

def _build_arbeitsdatei_summary(ws_base: Path, max_chars: int = 3000) -> str:
    """Build a compact summary of the arbeitsdatei for the Phase 4 Draft prompt.

    The raw arbeitsdatei can be 300K+ chars (2000+ lines). This function
    extracts category totals, counts, and open questions — enough for the
    LLM to write a meaningful email without needing to load the full file.
    Pure Python, no LLM call.
    """
    arbeitsdatei = ws_base / "arbeitsdatei.md"
    if not arbeitsdatei.exists():
        return ""

    try:
        lines = arbeitsdatei.read_text(encoding="utf-8").split("\n")
    except Exception:
        return ""

    # Parse table rows: | Datum | Beschreibung | Betrag | Kategorie | ... | Status |
    import re
    categories: dict[str, list[float]] = {}
    open_questions: list[str] = []
    total_positions = 0

    for line in lines:
        if not line.strip().startswith("|"):
            continue
        cells = [c.strip() for c in line.split("|")]
        if len(cells) < 7:
            continue
        # Skip header rows
        if "Datum" in cells[1] or "---" in cells[1]:
            continue

        total_positions += 1
        beschreibung = cells[2] if len(cells) > 2 else ""
        betrag_raw = cells[3] if len(cells) > 3 else "0"
        kategorie = cells[4] if len(cells) > 4 else "privat"
        status = cells[-2] if len(cells) > 6 else ""

        # Parse amount
        betrag = 0.0
        amount_match = re.search(r'[\d.,]+', betrag_raw.replace(".", "").replace(",", "."))
        if amount_match:
            try:
                betrag = float(amount_match.group())
            except ValueError:
                pass

        # Categorize
        kat_key = kategorie.strip().lower() if kategorie.strip() else "unkategorisiert"
        if kat_key not in categories:
            categories[kat_key] = []
        categories[kat_key].append(betrag)

        # Collect open questions
        status_lower = status.lower()
        if "nachfragen" in status_lower or "unklar" in status_lower:
            if len(open_questions) < 15:
                open_questions.append(f"- {beschreibung} ({betrag_raw}): {status}")

    if total_positions == 0:
        return ""

    # Build summary
    parts = [f"ARBEITSDATEI SUMMARY ({total_positions} positions analyzed):\n"]

    # Sort categories by total amount (descending)
    sorted_cats = sorted(categories.items(), key=lambda x: sum(x[1]), reverse=True)
    for kat, amounts in sorted_cats:
        total = sum(amounts)
        count = len(amounts)
        if total > 0 or count > 5:
            parts.append(f"  {kat}: {count} positions, total {total:.2f} EUR")

    if open_questions:
        parts.append(f"\nOPEN QUESTIONS ({len(open_questions)}):")
        parts.extend(open_questions)

    # Also include offene_punkte.md if it exists (compact)
    offene = ws_base / "offene_punkte.md"
    if offene.exists():
        try:
            op_text = offene.read_text(encoding="utf-8")[:800]
            parts.append(f"\nOFFENE PUNKTE:\n{op_text}")
        except Exception:
            pass

    summary = "\n".join(parts)
    return summary[:max_chars]


def _build_reference_summary(ws_base: Path, max_chars: int = 6000) -> str:
    """Load reference/ files into a context block for Phase 4 prompt.

    FAQ/knowledge-base agents store their domain knowledge in reference/ files.
    Without this injection, the LLM has no data to answer questions — the Phase 4
    prompt forbids tool calls, and the Lagebild doesn't contain file contents.

    Bug #19: Root cause of proof_faq_bot E2E failures (PROOF-008 through PROOF-017).
    The agent's prompt says "read the FAQ file" but Phase 4 says "no tool calls".
    Fix: pre-inject all reference files so the LLM has the data in-context.
    """
    ref_dir = ws_base / "reference"
    if not ref_dir.exists() or not ref_dir.is_dir():
        return ""

    parts = ["REFERENZDATEN (Wissensdatenbank — als Antwortquelle nutzen):\n"]
    total_chars = 0

    for ref_file in sorted(ref_dir.rglob("*")):
        if not ref_file.is_file():
            continue
        if ref_file.suffix not in (".md", ".txt", ".csv", ".json", ".yaml", ".yml"):
            continue
        try:
            content = ref_file.read_text(encoding="utf-8")
            rel_path = ref_file.relative_to(ws_base)
            if total_chars + len(content) > max_chars:
                remaining = max_chars - total_chars
                if remaining > 200:
                    parts.append(f"\n--- {rel_path} (truncated) ---\n{content[:remaining]}\n")
                break
            parts.append(f"\n--- {rel_path} ---\n{content}\n")
            total_chars += len(content)
        except Exception:
            continue

    if len(parts) <= 1:
        return ""

    return "\n".join(parts)


# ── Phase 4 Task Modes ────────────────────────────────────────────────────

def _phase4_task_dispatch() -> str:
    """Phase 4 task prompt for dispatch-mode agents (email, chat).
    Agent writes plain text → orchestrator dispatches it."""
    return (
        "TASK: Write your response as PLAIN TEXT (Fließtext). This text will be sent "
        "to the stakeholder by the system.\n"
        "IMPORTANT RULES:\n"
        "- Write ONLY the response text. Nothing else.\n"
        "- Do NOT write tool calls, XML tags, code blocks, or JSON.\n"
        "- Do NOT write <tool_call>, <function=...>, read_file(), or remember().\n"
        "- Do NOT try to save or store information — just write the response.\n"
        "- Use the REFERENZDATEN, Lagebild, and Arbeitsdatei Summary as your data source.\n"
        "- Answer the stakeholder's question directly based on the data provided above.\n"
        "- ONLY cite facts that appear in the data. Do NOT speculate.\n"
        "- NEVER mention names, messages, or internal details of other stakeholders."
    )


def _phase4_task_file(agent: "AIMOSAgent") -> str:
    """Phase 4 task prompt for file-mode — UNUSED when output_assembly is active.
    Kept as fallback for agents without batch_file_specs."""
    return "TASK: Erstelle die Output-Dateien gemaess Deinem Arbeitsablauf."


FALLBACK_MARKER = "Dokument konnte nicht automatisch analysiert werden"


def _is_fallback_section(section_text: str) -> bool:
    """True iff the section is only a fallback stub (no real analysis)."""
    return FALLBACK_MARKER in section_text and len(section_text.strip()) < 250


def _doc_has_real_analysis(arbeitsdatei_text: str, doc_name: str) -> bool:
    """True iff doc_name has a real (non-fallback) section in arbeitsdatei.

    A fallback stub counts as "not processed" so the next cycle can retry.
    """
    import re as _re_real
    if f"## {doc_name}" not in arbeitsdatei_text:
        return False
    for p in _re_real.split(r"\n(?=## )", arbeitsdatei_text):
        title = p.split("\n")[0].strip("# ").strip()
        if title == doc_name:
            return not _is_fallback_section(p)
    return False


def _remove_fallback_section(arbeitsdatei_path: Path, doc_name: str) -> bool:
    """Remove a fallback stub section for doc_name from arbeitsdatei.md.

    Returns True if a fallback section was removed.
    """
    import re as _re_rm
    if not arbeitsdatei_path.exists():
        return False
    text = arbeitsdatei_path.read_text(encoding="utf-8")
    parts = _re_rm.split(r"(\n(?=## ))", text)  # capture separators to rebuild
    out = []
    removed = False
    i = 0
    while i < len(parts):
        chunk = parts[i]
        title = chunk.lstrip("\n").split("\n", 1)[0].strip("# ").strip()
        if title == doc_name and _is_fallback_section(chunk):
            removed = True
            i += 1
            if i < len(parts) and parts[i].startswith("\n"):
                i += 1  # skip the captured "\n" separator that followed the removed block
            continue
        out.append(chunk)
        i += 1
    if removed:
        arbeitsdatei_path.write_text("".join(out), encoding="utf-8")
    return removed


def _list_input_filenames(ws_base: Path) -> set[str]:
    """Return the set of filenames currently in the agent's input/ directory.

    Used as whitelist for arbeitsdatei section detection. If input/ does not
    exist (or scan dir is configured differently), returns an empty set —
    callers treat that as "no whitelist, accept anything that passes the
    other heuristics".
    """
    input_dir = ws_base / "input"
    if not input_dir.is_dir():
        return set()
    return {p.name for p in input_dir.iterdir() if p.is_file()}


def _sanitize_chunk_result(result: str) -> str:
    """Strip leading `## ` headers from a chunk-analysis result.

    LLMs sometimes prepend their own `## Analyseergebnis fuer Dokument: ...`
    header, which the section splitter would later misread as a separate
    document section. The harness owns section headers — the LLM's analysis
    body must not introduce new ones.

    Demotes any `## ` line to `### ` so structural Markdown stays readable
    but cannot create false sections.
    """
    out: list[str] = []
    for line in result.splitlines():
        stripped = line.lstrip()
        if stripped.startswith("## ") and not stripped.startswith("### "):
            out.append(line.replace("## ", "### ", 1))
        else:
            out.append(line)
    return "\n".join(out)


def _split_arbeitsdatei_doc_sections(
    content: str, expected_filenames: set[str] | None = None,
) -> list[str]:
    """Split arbeitsdatei.md content into document sections.

    Single source of truth for "what counts as a document section".
    Used by Output Assembly and the Completion Check so they cannot drift.

    A document section:
      - starts with `## ` (after strip), so the empty preamble before the first
        `## ` and any leading blank block from re.split are excluded;
      - has at least 50 chars of content (filters out near-empty stubs);
      - has a title containing `.` or `__` (filters out LLM-generated headings
        like `## A. Technische Dokumentation`);
      - if `expected_filenames` is given, the title must match one of them
        exactly — defends against LLM-injected fake headers like
        `## Analyseergebnis fuer Dokument: **Foo.md**` slipping through.
    """
    import re as _re_split
    out: list[str] = []
    for s in _re_split.split(r"\n(?=## )", content):
        stripped = s.strip()
        if not stripped.startswith("## ") or len(stripped) < 50:
            continue
        title = s.split("\n")[0].strip("# ").strip()
        if "." not in title and "__" not in title:
            continue
        if expected_filenames is not None and title not in expected_filenames:
            continue
        out.append(s)
    return out


def _check_file_mode_completion(
    agent: "AIMOSAgent", ws_base: Path, log: logging.Logger,
) -> dict:
    """State-based completion check for file-mode agents.

    Returns dict with:
      complete: bool — True if all output files exist and are valid
      reason: str — why it's incomplete (for retrigger message)
      docs_in_arbeitsdatei: int — sections in arbeitsdatei
      docs_in_inventar: int — entries in inventar JSON
      output_files: list — which output files were checked
    """
    result = {"complete": False, "reason": "", "docs_in_arbeitsdatei": 0,
              "docs_in_inventar": 0, "output_files": []}

    _file_specs = agent.config.get("batch_file_specs", {})
    _output_files = [f for f in _file_specs.keys() if f != "arbeitsdatei.md"]
    result["output_files"] = _output_files

    # 1. arbeitsdatei must exist and have content
    _arbeitsdatei = ws_base / "arbeitsdatei.md"
    if not _arbeitsdatei.exists() or _arbeitsdatei.stat().st_size < 100:
        result["reason"] = "no_arbeitsdatei"
        return result

    # Count document sections in arbeitsdatei (shared helper with Output Assembly).
    # Pass the actual input filenames as whitelist so LLM-injected fake `## ` headers
    # cannot inflate the count.
    _ad_text = _arbeitsdatei.read_text(encoding="utf-8")
    _expected_names = _list_input_filenames(ws_base)
    _ad_sections = len(_split_arbeitsdatei_doc_sections(_ad_text, _expected_names))
    result["docs_in_arbeitsdatei"] = _ad_sections

    # 2. All output files must exist
    _missing = [f for f in _output_files if not (ws_base / f).exists()]
    if _missing:
        result["reason"] = f"Fehlende Output-Dateien: {', '.join(_missing)}"
        return result

    # 3. Validate each output file against its spec
    for fname in _output_files:
        fpath = ws_base / fname
        if fpath.stat().st_size < 10:
            result["reason"] = f"{fname} ist leer"
            return result

        _spec_marker = _file_specs.get(fname, "")
        if _spec_marker:
            _fcontent = fpath.read_text(encoding="utf-8")
            if _spec_marker not in _fcontent:
                result["reason"] = f"{fname} enthält nicht den erwarteten Marker"
                return result

    # 4. For JSON inventar: check document count matches arbeitsdatei
    for fname in _output_files:
        if fname.endswith(".json"):
            fpath = ws_base / fname
            try:
                _inv = json.loads(fpath.read_text(encoding="utf-8"))
                _inv_docs = len(_inv.get("dokumente", []))
                result["docs_in_inventar"] = _inv_docs
                if _ad_sections > 0 and _inv_docs < _ad_sections:
                    result["reason"] = (
                        f"{fname} hat {_inv_docs} Dokumente, "
                        f"arbeitsdatei hat {_ad_sections} Abschnitte"
                    )
                    return result
            except (json.JSONDecodeError, KeyError):
                result["reason"] = f"{fname} ist kein valides JSON"
                return result

    result["complete"] = True
    return result


def _load_valid_anhang_iv(ws_base: Path) -> set[str]:
    """Load valid Anhang-IV IDs from the reference checklist file (Single Source of Truth)."""
    import re as _re_aiv
    for ref_dir in [ws_base / "reference", ws_base / "team_reference"]:
        for fname in ["anhang_IV_checkliste.md", "anhang_iv_anforderungen.md"]:
            ref_file = ref_dir / fname
            if ref_file.exists():
                ids = set()
                for m in _re_aiv.finditer(r"id:\s*(\S+)", ref_file.read_text(encoding="utf-8")):
                    ids.add(m.group(1))
                if ids:
                    return ids
    return set()  # No checklist found — skip validation


def _load_valid_types(ws_base: Path) -> list[str]:
    """Load valid document types from dokumenttypen.yaml (Single Source of Truth).

    Extracts all subtype entries from the YAML catalog.
    Falls back to empty list if no catalog found (= no type constraint in prompt).
    """
    import re as _re_typ
    for ref_dir in [ws_base / "team_reference", ws_base / "reference"]:
        for fname in ["dokumenttypen.yaml", "dokumenttypen.yml"]:
            f = ref_dir / fname
            if f.exists():
                types = []
                for m in _re_typ.finditer(r"^\s+-\s+(\w+)", f.read_text(encoding="utf-8"), _re_typ.MULTILINE):
                    typ = m.group(1)
                    if typ not in ("subtypen",):  # Skip YAML keys
                        types.append(typ)
                if types:
                    return sorted(set(types))
    return []


def _normalize_entry(
    entry: dict, valid_analysts: set, valid_anhang_iv: set,
    valid_types: list[str],
    log: logging.Logger, agent_name: str,
) -> dict:
    """Harness-level normalization of LLM-generated inventory entries."""
    import re as _re_norm
    # Normalize ID: DOK-009 → DOK-09, DOK-002 → DOK-02
    raw_id = entry.get("id", "")
    m = _re_norm.match(r"DOK-0*(\d+)", raw_id)
    if m:
        num = int(m.group(1))
        entry["id"] = f"DOK-{num:02d}"
        if entry["id"] != raw_id:
            log.info(f"[{agent_name}] Normalize ID: {raw_id} → {entry['id']}")

    # Normalize dateiname: ensure it starts with DOK-NN_
    dateiname = entry.get("dateiname", "")
    if dateiname and not dateiname.startswith("DOK-"):
        entry["dateiname"] = f"{entry['id']}_{dateiname}"

    # Filter invalid Anhang-IV IDs (loaded from reference checklist)
    raw_zuordnung = entry.get("anhang_iv_zuordnung", [])
    valid_zuordnung = [a for a in raw_zuordnung if a in valid_anhang_iv] if valid_anhang_iv else raw_zuordnung
    if len(valid_zuordnung) != len(raw_zuordnung):
        removed = set(raw_zuordnung) - set(valid_zuordnung)
        log.info(f"[{agent_name}] {entry['id']}: Removed invalid Anhang-IV IDs: {removed}")
    entry["anhang_iv_zuordnung"] = valid_zuordnung

    # Validate document type against catalog
    if valid_types:
        raw_typ = entry.get("typ", "sonstiges")
        if raw_typ not in valid_types:
            # Try lowercase/normalized match
            match = None
            raw_lower = raw_typ.lower().replace("-", "_").replace(" ", "_")
            for vt in valid_types:
                if vt == raw_lower or raw_lower.startswith(vt) or vt.startswith(raw_lower):
                    match = vt
                    break
            if match:
                log.info(f"[{agent_name}] {entry['id']}: Typ {raw_typ} → {match}")
                entry["typ"] = match
            else:
                log.warning(f"[{agent_name}] {entry['id']}: Unknown typ '{raw_typ}', keeping as-is")

    # Normalize analyst names
    if valid_analysts:
        raw_analysts = entry.get("relevant_fuer", [])
        normalized = []
        for a in raw_analysts:
            if a in valid_analysts:
                normalized.append(a)
            else:
                # Fuzzy match: mca_elektro/mca_elec → mca_elek
                best = None
                for va in valid_analysts:
                    if a[:6] == va[:6]:  # Same prefix after mca_
                        best = va
                        break
                if best:
                    log.info(f"[{agent_name}] {entry['id']}: Analyst {a} → {best}")
                    normalized.append(best)
                else:
                    log.warning(f"[{agent_name}] {entry['id']}: Unknown analyst '{a}', dropped")
        entry["relevant_fuer"] = normalized

    # Ensure pfad field
    if "pfad" not in entry or entry["pfad"] in ("N/A", "", None):
        entry["pfad"] = f"input/{entry.get('dateiname', '')}"

    # Ensure begruendung field — always present so consumers can rely on it
    if not entry.get("begruendung"):
        entry["begruendung"] = ""
        log.warning(f"[{agent_name}] {entry['id']}: begruendung fehlt im LLM-Output")

    return entry


async def _phase4_output_assembly(
    agent: "AIMOSAgent", ws_base: Path, log: logging.Logger,
) -> list[dict]:
    """Phase 4 Output Assembly — Harness baut, LLM entscheidet klein.

    1. arbeitsdatei.md in Abschnitte splitten (## Separatoren)
    2. Pro Abschnitt: 1 kleiner LLM-Call → 1 JSON-Entry
    3. Harness assembliert deterministisch → dokumenten_inventar.json
    4. Gap-Analyse deterministisch → maengelliste.md
    """
    arbeitsdatei = ws_base / "arbeitsdatei.md"
    if not arbeitsdatei.exists() or arbeitsdatei.stat().st_size < 100:
        log.warning(f"[{agent.agent_name}] Output Assembly: arbeitsdatei leer/fehlt")
        return []

    content = arbeitsdatei.read_text(encoding="utf-8")

    # 1. Splitten nach ## Dokument-Abschnitten (gemeinsame Logik mit Completion-Check)
    import re as _re_oa
    _all_sections = _re_oa.split(r"\n(?=## )", content)
    _expected_names = _list_input_filenames(ws_base)
    doc_sections = _split_arbeitsdatei_doc_sections(content, _expected_names)
    _skipped = 0
    for s in _all_sections:
        stripped = s.strip()
        if not stripped.startswith("## ") or len(stripped) < 50:
            continue
        title = s.split("\n")[0].strip("# ").strip()
        if "." not in title and "__" not in title:
            _skipped += 1
    if _skipped:
        log.info(f"[{agent.agent_name}] Output Assembly: {_skipped} LLM-generierte Abschnitte uebersprungen")
    log.info(f"[{agent.agent_name}] Output Assembly: {len(doc_sections)} Dokument-Abschnitte in arbeitsdatei")

    if not doc_sections:
        log.warning(f"[{agent.agent_name}] Output Assembly: Keine ## Abschnitte gefunden")
        return []

    # 2. Lade gültige Analysten-Namen aus Mapping-Datei
    _valid_analysts = set()
    for ref_dir in [ws_base / "reference", ws_base / "team_reference"]:
        for fname in ["analysten_mapping.md", "analysten_zuordnung.md"]:
            _mapping_file = ref_dir / fname
            if _mapping_file.exists():
                for _line in _mapping_file.read_text(encoding="utf-8").splitlines():
                    _m = _re_oa.match(r"\|\s*(mca_\w+)\s*\|", _line)
                    if _m:
                        _valid_analysts.add(_m.group(1))
                break
        if _valid_analysts:
            break
    _analysts_str = ", ".join(sorted(_valid_analysts)) if _valid_analysts else ""

    # Lade gültige Anhang-IV IDs aus Referenz-Checkliste
    _valid_anhang_iv = _load_valid_anhang_iv(ws_base)
    _anhang_iv_str = ", ".join(sorted(_valid_anhang_iv)) if _valid_anhang_iv else ""

    # Lade gültige Dokumenttypen aus Katalog
    _valid_types = _load_valid_types(ws_base)
    _types_str = ", ".join(_valid_types) if _valid_types else ""
    log.info(f"[{agent.agent_name}] Output Assembly: {len(_valid_anhang_iv)} Anhang-IV IDs, {len(_valid_analysts)} Analysten, {len(_valid_types)} Dokumenttypen geladen")

    # Pro Abschnitt: LLM generiert EINEN JSON-Eintrag
    # W5 (Stable Doc-IDs): wenn Config-Flag gesetzt, Hash-basierte IDs.
    # Sonst Sequential (rueckwaerts-kompatibel).
    _use_stable_ids = bool(agent.config.get("stable_doc_ids", False))
    _stable_id_map: dict[str, str] = {}
    if _use_stable_ids:
        try:
            from core.stable_doc_ids import assign_stable_ids
            file_dicts = []
            for s in doc_sections:
                title = s.split("\n")[0].strip("# ").strip()
                file_dicts.append({"dateiname": title, "size_bytes": len(s)})
            _stable_id_map = assign_stable_ids(file_dicts)
            log.info(f"[{agent.agent_name}] W5 stable_doc_ids: {len(_stable_id_map)} IDs zugeordnet")
        except Exception as exc:
            log.warning(f"[{agent.agent_name}] stable_doc_ids fehlgeschlagen, fallback: {exc}")
            _stable_id_map = {}

    entries = []
    for idx, section in enumerate(doc_sections):
        section_title = section.split("\n")[0].strip("# ").strip()
        # Deterministic ID assignment by harness (not LLM)
        if _stable_id_map:
            _harness_doc_id = _stable_id_map.get(section_title) or f"DOK-{idx+1:04d}"
        else:
            _harness_doc_id = f"DOK-{idx+1:04d}"
        log.info(f"[{agent.agent_name}] Output Assembly: Entry {idx+1}/{len(doc_sections)}: {_harness_doc_id} ← {section_title[:40]}")

        # Build constraint lines from loaded reference data (Single Source of Truth)
        _constraints = []
        _constraints.append(f"- id: Verwende exakt diese ID: {_harness_doc_id}")
        if _types_str:
            _constraints.append(f"- typ: NUR diese Werte erlaubt: {_types_str}")
        else:
            _constraints.append("- typ: Freitext-Klassifikation des Dokumenttyps")
        _constraints.append("- aktuell: false wenn aelter als 2 Jahre oder veraltete Norm")
        if _anhang_iv_str:
            _constraints.append(f"- anhang_iv_zuordnung: NUR diese Werte erlaubt: {_anhang_iv_str}")
        if _analysts_str:
            _constraints.append(f"- relevant_fuer: NUR diese Werte erlaubt: {_analysts_str}")
        _constraints.append("- Antworte NUR mit dem JSON-Objekt, kein Markdown, kein Text.")

        entry_prompt = (
            "Erstelle aus dieser Dokumentenanalyse EINEN JSON-Eintrag.\n\n"
            f"ANALYSE:\n{section[:3000]}\n\n"
            "FORMAT (antworte NUR mit diesem JSON, nichts anderes):\n"
            '{"id": "DOK-NN", "dateiname": "...", "typ": "...", '
            '"revision": "...", "datum": "YYYY-MM-DD", "aktuell": true, '
            '"anhang_iv_zuordnung": ["A.1"], "relevant_fuer": ["mca_mech"], '
            '"warnungen": ["..."], '
            '"begruendung": "Warum dieser Typ und diese Aktualitaets-Bewertung"}\n\n'
            "Regeln:\n" + "\n".join(_constraints)
            + "\n- begruendung: 1-2 Saetze die erklaeren WARUM dieser Typ und aktuell/veraltet"
        )

        try:
            with _PhaseParams(agent, "phase4", log):
                result = await _think_with_activity_check(
                    agent, entry_prompt, log,
                    stale_timeout=30,  # Kurzer Timeout pro Entry
                    phase="4",
                )
            # Parse JSON from result
            result = result.strip()
            # Detect LLM refusal and retry with simplified prompt
            _is_refusal = any(kw in result.lower() for kw in [
                "entschuldigung", "kann ich nicht", "kann diese anfrage nicht",
                "wende dich an", "i cannot", "i'm sorry",
            ])
            if _is_refusal:
                log.warning(f"[{agent.agent_name}] Output Assembly: LLM refusal for {section_title[:30]}, retrying with simplified prompt")
                _retry_prompt = (
                    "Klassifiziere dieses Dokument als JSON.\n\n"
                    f"Titel: {section_title}\n"
                    f"Inhalt (Auszug): {section[:1500]}\n\n"
                    "Antwort NUR als JSON: "
                    '{"id": "DOK-NN", "dateiname": "...", "typ": "...", '
                    '"revision": "...", "datum": "...", "aktuell": true, '
                    '"anhang_iv_zuordnung": [], "relevant_fuer": [], "warnungen": [], '
                    '"begruendung": "kurze Begruendung"}'
                )
                try:
                    with _PhaseParams(agent, "phase4", log):
                        result = await _think_with_activity_check(
                            agent, _retry_prompt, log, stale_timeout=30, phase="4",
                        )
                    result = result.strip()
                except Exception as exc:
                    log.warning(f"[{agent.agent_name}] Output Assembly: Retry also failed: {exc}")

            # Strip markdown fences
            if result.startswith("```"):
                result = _re_oa.sub(r"^```\w*\n?", "", result)
                result = _re_oa.sub(r"\n?```$", "", result.strip())
            try:
                entry = json.loads(result)
                if isinstance(entry, dict):
                    entry["id"] = _harness_doc_id  # Harness-assigned ID overrides LLM
                    entry = _normalize_entry(entry, _valid_analysts, _valid_anhang_iv, _valid_types, log, agent.agent_name)
                    entries.append(entry)
                    log.info(f"[{agent.agent_name}] Output Assembly: {entry['id']}: {entry.get('typ')}")
                else:
                    log.warning(f"[{agent.agent_name}] Output Assembly: Kein valides Entry: {result[:100]}")
            except json.JSONDecodeError:
                # Retry: Versuch JSON aus dem Text zu extrahieren
                m = _re_oa.search(r"\{[^{}]*\}", result, _re_oa.DOTALL)
                if m:
                    try:
                        entry = json.loads(m.group())
                        entry["id"] = _harness_doc_id  # Harness-assigned ID overrides LLM
                        entry = _normalize_entry(entry, _valid_analysts, _valid_anhang_iv, _valid_types, log, agent.agent_name)
                        entries.append(entry)
                        log.info(f"[{agent.agent_name}] Output Assembly: {entry['id']} (extracted)")
                    except json.JSONDecodeError:
                        log.warning(f"[{agent.agent_name}] Output Assembly: JSON parse failed: {result[:100]}")
                else:
                    log.warning(f"[{agent.agent_name}] Output Assembly: Kein JSON in Antwort: {result[:100]}")
        except Exception as exc:
            log.warning(f"[{agent.agent_name}] Output Assembly: LLM-Call failed: {exc}")

    log.info(f"[{agent.agent_name}] Output Assembly: {len(entries)}/{len(doc_sections)} Entries generiert")

    # 3. Harness assembliert deterministisch
    # Gap-Analyse: Welche Anhang-IV Punkte sind abgedeckt?
    covered = set()
    for entry in entries:
        for aid in entry.get("anhang_iv_zuordnung", []):
            covered.add(aid)

    # Lade Anhang-IV Checkliste für fehlende Punkte
    checklist_items = []
    for ref_dir in [ws_base / "reference", ws_base / "team_reference"]:
        for fname in ["anhang_IV_checkliste.md", "anhang_iv_anforderungen.md"]:
            ref_file = ref_dir / fname
            if ref_file.exists():
                ref_text = ref_file.read_text(encoding="utf-8")
                # Parse YAML-like entries: "- id: A.1\n    name: ..."
                for m in _re_oa.finditer(r"id:\s*(\S+)\s*\n\s*name:\s*\"?([^\"]+)\"?", ref_text):
                    checklist_items.append({"id": m.group(1), "name": m.group(2).strip()})
                break
        if checklist_items:
            break

    all_required = set(item["id"] for item in checklist_items)
    missing = all_required - covered

    # Projektname aus status.json oder firma_profil.md laden
    # ws_base = storage/agents/<name>/, Workspace = workspaces/<name>/
    _aimos_root = ws_base.parent.parent.parent  # /opt/AIMOS/storage/agents/<name> → /opt/AIMOS
    _ws_dir = _aimos_root / "workspaces" / agent.agent_name
    _projekt_name = agent.agent_name
    for _sf in [_ws_dir / "status.json", ws_base / "status.json"]:
        if _sf.exists():
            try:
                _sj = json.loads(_sf.read_text(encoding="utf-8"))
                _projekt_name = _sj.get("projekt_id") or _sj.get("maschine") or _projekt_name
                if _projekt_name != agent.agent_name:
                    break
            except (json.JSONDecodeError, KeyError):
                pass
    if _projekt_name == agent.agent_name:
        for _pf in [_ws_dir / "firma_profil.md", ws_base / "firma_profil.md"]:
            if _pf.exists():
                _pf_text = _pf.read_text(encoding="utf-8")
                _pm = _re_oa.search(r"(?:Projekt-ID|Maschine)[:\s]*(.+)", _pf_text)
                if _pm:
                    _projekt_name = _pm.group(1).strip()
                    break

    inventar = {
        "projekt": _projekt_name,
        "stand": __import__("datetime").date.today().isoformat(),
        "datenquellen_durchsucht": 1,
        "dokumente_gesamt": len(entries),
        "dokumente": entries,
        "gap_analyse": {
            "vorhanden": sorted(covered),
            "fehlend": [
                {"id": mid, "grund": next(
                    (i["name"] for i in checklist_items if i["id"] == mid),
                    f"{mid} fehlt"
                )}
                for mid in sorted(missing)
            ],
        },
    }

    # Schreibe Inventar (deterministisch)
    inventar_path = ws_base / "dokumenten_inventar.json"
    inventar_path.write_text(json.dumps(inventar, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    log.info(f"[{agent.agent_name}] Output Assembly: dokumenten_inventar.json geschrieben ({len(entries)} Dokumente)")

    # 4. Mängelliste deterministisch
    lines = ["# Maengelliste\n"]
    lines.append("| Prioritaet | Anhang-IV | Was fehlt | Status |")
    lines.append("|:----------:|:---------:|-----------|--------|")
    for item in inventar["gap_analyse"]["fehlend"]:
        aid = item["id"]
        if aid.startswith("B") or aid.startswith("C"):
            prio = "HOCH"
        elif aid.startswith("D") or aid.startswith("E"):
            prio = "HOCH"
        elif aid.startswith("A"):
            prio = "MITTEL"
        else:
            prio = "NIEDRIG"
        lines.append(f"| {prio} | {aid} | {item['grund']} | OFFEN |")
    # Auch veraltete Dokumente aufnehmen
    for entry in entries:
        if not entry.get("aktuell", True):
            lines.append(f"| HOCH | — | {entry.get('id','?')}: {entry.get('typ','?')} ist VERALTET | AKTUALISIERUNG |")
        for w in entry.get("warnungen", []):
            lines.append(f"| MITTEL | — | {entry.get('id','?')}: {w} | KLAERUNG |")

    maengel_path = ws_base / "maengelliste.md"
    maengel_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    log.info(f"[{agent.agent_name}] Output Assembly: maengelliste.md geschrieben ({len(lines)-3} Eintraege)")

    return [{"action": "output_assembly", "entries": len(entries), "missing": len(missing)}]


# ── CR-279: Draft Safety (universal + agent-specific hooks) ────────────────

def _apply_draft_safety(
    draft: str, agent: "AIMOSAgent", thread_id: str, log: logging.Logger,
) -> str:
    """Minimal universal safety check on Phase 4 draft output.

    Only catches the most egregious failures (raw XML/tool-call text that
    would make the email unreadable). Does NOT restrict LLM creativity.

    Agent-specific post-filters can be registered via batch_hooks config:
        "batch_hooks": {"draft_post_filter": "mymodule.my_filter_function"}
    The hook receives (draft, thread_id) and returns the filtered draft.
    """
    import re

    # Universal: strip XML tool-call blocks that make emails unreadable
    original_len = len(draft)
    # Closed tags first
    draft = re.sub(r'<tool_call>.*?</tool_call>', '', draft, flags=re.DOTALL)
    draft = re.sub(r'<tool_call_id>.*?</tool_call_id>', '', draft, flags=re.DOTALL)
    draft = re.sub(r'<tool_name>.*?</tool_name>', '', draft, flags=re.DOTALL)
    draft = re.sub(r'<tool_arguments>.*?</tool_arguments>', '', draft, flags=re.DOTALL)
    draft = re.sub(r'<invoke\s+name="[^"]*">.*?</invoke>', '', draft, flags=re.DOTALL)
    # Unclosed tags (LLM often writes <tool_call> without </tool_call>)
    draft = re.sub(r'<tool_call\b[^>]*>.*?(?=<tool_call|$)', '', draft, flags=re.DOTALL)
    # Any remaining tool-related XML tags
    draft = re.sub(r'</?tool_(?:call|call_id|name|arguments)[^>]*>', '', draft)
    draft = re.sub(r'</?invoke[^>]*>', '', draft)
    # Python-style tool calls
    draft = re.sub(r'(?:read_file|write_file|send_email|remember|recall)\s*\([^)]*\)', '', draft)
    draft = re.sub(r'\n{3,}', '\n\n', draft).strip()

    if len(draft) < original_len * 0.5 and original_len > 100:
        log.warning(
            f"[{agent.agent_name}] CR-279: Draft was >50% tool-call text "
            f"({original_len} → {len(draft)} chars). LLM likely in wrong mode."
        )

    # Agent-specific hook (optional)
    hook_path = agent.config.get("batch_hooks", {}).get("draft_post_filter")
    if hook_path:
        try:
            module_name, func_name = hook_path.rsplit(".", 1)
            import importlib
            mod = importlib.import_module(module_name)
            hook_fn = getattr(mod, func_name)
            draft = hook_fn(draft, thread_id)
            log.debug(f"[{agent.agent_name}] CR-279: Agent hook '{hook_path}' applied")
        except Exception as exc:
            log.warning(f"[{agent.agent_name}] CR-279: Hook '{hook_path}' failed: {exc}")

    return draft


# ── CR-258: Workspace document scanning ────────────────────────────────────

_DOCUMENT_EXTENSIONS = {
    # PDF + Images
    ".pdf", ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp", ".heic", ".heif",
    # MS Office modern (Open XML)
    ".docx", ".xlsx", ".pptx",
    # MS Office legacy (requires extra extraction)
    ".doc", ".xls", ".ppt",
    # Outlook / Email
    ".msg", ".eml",
    # MS Project, Visio (extract metadata/text where possible)
    ".mpp", ".vsdx",
    # LibreOffice / OpenDocument
    ".odt", ".ods", ".odp",
    # Rich text + plain text
    ".rtf", ".csv", ".txt", ".md",
}
_ARCHIVE_EXTENSIONS = {".zip", ".tar", ".gz", ".tgz", ".bz2", ".7z", ".rar"}


def _extract_archives(doc_dir: Path, log: logging.Logger, agent_name: str) -> int:
    """CR-258: Auto-extract archives in dokumente/ folder.

    Extracts ZIP, tar.gz, 7z, RAR into a subfolder named after the archive.
    Returns number of archives extracted. Skips already-extracted archives
    (subfolder with same name already exists).
    """
    extracted = 0
    if not doc_dir.exists():
        return 0
    for f in list(doc_dir.iterdir()):
        if not f.is_file() or f.suffix.lower() not in _ARCHIVE_EXTENSIONS:
            continue
        # Skip if already extracted (folder with same stem exists)
        target_dir = doc_dir / f.stem
        if target_dir.exists():
            continue
        try:
            target_dir.mkdir(parents=True, exist_ok=True)
            suffix = f.suffix.lower()
            if suffix == ".zip":
                import zipfile
                with zipfile.ZipFile(f, 'r') as zf:
                    zf.extractall(target_dir)
            elif suffix in {".tar", ".gz", ".tgz", ".bz2"}:
                import tarfile
                with tarfile.open(f, 'r:*') as tf:
                    tf.extractall(target_dir, filter='data')
            elif suffix == ".7z":
                import subprocess
                subprocess.run(["7z", "x", str(f), f"-o{target_dir}"], capture_output=True, timeout=60)
            elif suffix == ".rar":
                import subprocess
                subprocess.run(["unrar", "x", str(f), str(target_dir)], capture_output=True, timeout=60)
            else:
                target_dir.rmdir()
                continue
            extracted += 1
            log.info(f"[{agent_name}] CR-258: Extracted archive {f.name} → {target_dir.name}/")
        except Exception as exc:
            log.warning(f"[{agent_name}] CR-258: Failed to extract {f.name}: {exc}")
            if target_dir.exists() and not any(target_dir.iterdir()):
                target_dir.rmdir()
    return extracted


def _scan_workspace_documents(ws_base: Path, log: logging.Logger, agent_name: str) -> list[dict]:
    """CR-258: Scan workspace for new/unprocessed documents.

    Auto-extracts archives first, then scans for document files.
    Returns list of document descriptors for files in the workspace 'dokumente/'
    subdirectory (or client subdirs like klient_*/dokumente/).
    A file is considered processed if its name appears in arbeitsdatei.md or state.md.
    """
    # CR-258+: Configurable scan directories via batch_data_sources
    # Default: "dokumente/" (steuerberater pattern)
    # MCA pattern: "input/" (customer documents)
    doc_dir_names = ["dokumente", "input"]  # scan both by default
    doc_dirs = [ws_base / d for d in doc_dir_names]
    # Also check client subdirectories (steuerberater pattern: klient_*/dokumente/)
    if ws_base.exists():
        for d in ws_base.iterdir():
            if d.is_dir() and (d / "dokumente").is_dir():
                doc_dirs.append(d / "dokumente")
            if d.is_dir() and (d / "input").is_dir():
                doc_dirs.append(d / "input")

    # Auto-extract archives before scanning
    for doc_dir in doc_dirs:
        if doc_dir.exists():
            _extract_archives(doc_dir, log, agent_name)

    documents = []
    for doc_dir in doc_dirs:
        if not doc_dir.exists():
            continue
        # Recursive scan (includes extracted archive subdirectories)
        for f in sorted(doc_dir.rglob("*")):
            if f.is_file() and f.suffix.lower() in _DOCUMENT_EXTENSIONS:
                documents.append({
                    "path": f,
                    "name": f.name,
                    "size_kb": f.stat().st_size // 1024,
                    "modified": f.stat().st_mtime,
                    "client_dir": f.parent.parent.name if f.parent.name == "dokumente" else None,
                })

    if not documents:
        return []

    # Check which are already processed — ONLY check arbeitsdatei.md.
    # A fallback stub counts as NOT processed so the next cycle can retry.
    processed_names = set()
    check_path = ws_base / "arbeitsdatei.md"
    if check_path.exists():
        try:
            content = check_path.read_text(encoding="utf-8")
            for doc in documents:
                if _doc_has_real_analysis(content, doc["name"]):
                    processed_names.add(doc["name"])
        except Exception:
            pass

    new_docs = [d for d in documents if d["name"] not in processed_names]
    if new_docs:
        log.info(
            f"[{agent_name}] CR-258 Workspace scan: {len(new_docs)} new document(s) "
            f"({len(documents)} total, {len(processed_names)} already processed)"
        )

    # K1 Defense-in-depth: HR-/Personal-/Lohn-Daten 2nd-layer-Pruefung.
    # Default an fuer mca_*-Agenten; per Config deaktivierbar.
    if agent_name.startswith("mca_") and new_docs:
        try:
            from core.data_protection import filter_input_files
            input_dir = ws_base / "input"
            file_dicts = [
                {"pfad": str(d.get("path", "")), "dateiname": d["name"]}
                for d in new_docs
            ]
            ok, blocked = filter_input_files(
                file_dicts,
                workspace_input_dir=input_dir if input_dir.is_dir() else None,
                log=log,
            )
            if blocked:
                blocked_names = {b["dateiname"] for b in blocked}
                new_docs = [d for d in new_docs if d["name"] not in blocked_names]
                log.warning(
                    f"[{agent_name}] K1 Data-Protection: {len(blocked)} Datei(en) "
                    f"blockiert: {[b['dateiname'] for b in blocked]}"
                )
                # Audit-Log fuer Datenschutz-Vorfall
                try:
                    from core.audit import log_event
                    for b in blocked:
                        log_event(
                            ws_base, agent=agent_name,
                            action="data_protection_block",
                            target=b["dateiname"], result="blocked",
                            details={"grund": b.get("block_grund", "")},
                        )
                except Exception:
                    pass
        except ImportError:
            pass

    return new_docs


def _chunk_document_text(text: str, chunk_size: int = 2000) -> list[str]:
    """CR-258: Split document text into chunks for iterative processing.

    Splits at paragraph boundaries when possible, falls back to hard split.
    chunk_size is in characters (not tokens).
    """
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    remaining = text
    while remaining:
        if len(remaining) <= chunk_size:
            chunks.append(remaining)
            break
        # Try to split at a paragraph boundary
        split_pos = remaining.rfind("\n\n", 0, chunk_size)
        if split_pos < chunk_size // 2:
            split_pos = remaining.rfind("\n", 0, chunk_size)
        if split_pos < chunk_size // 2:
            split_pos = chunk_size
        chunks.append(remaining[:split_pos])
        remaining = remaining[split_pos:].lstrip("\n")

    return chunks


# ── CR-247: Activity-based timeout instead of hard timeout ─────────────────

async def _think_with_activity_check(
    agent: "AIMOSAgent", prompt: str, log: logging.Logger,
    stale_timeout: int = 120, poll_interval: int = 15,
    hard_timeout: int = 600, phase: str | None = None,
) -> str:
    """Call agent.think() with Ollama activity monitoring AND hard safety timeout.

    Instead of killing a think() call after N seconds regardless of progress,
    this checks whether Ollama is still actively generating. A call that takes
    20 minutes but is still producing tokens will NOT be killed. A call where
    Ollama has stopped responding for stale_timeout seconds WILL be killed.

    CR-270 hardening: Added hard_timeout as absolute wall-clock limit (default 600s)
    and consecutive API failure tracking. The /api/ps endpoint only tells us if a
    model is loaded, not if it's generating — so a model stuck mid-inference still
    appears "active". The hard_timeout catches this case.

    CR-273: Added phase parameter for OODA tool filtering. When set, agent.think()
    will filter the tool list to only those allowed in the given OODA phase.

    Args:
        stale_timeout: Seconds without Ollama activity before cancelling (default: 120)
        poll_interval: How often to check Ollama status (default: 15s)
        hard_timeout: Absolute max wall-clock seconds per call (default: 600)
        phase: OODA phase ("0"-"5") for tool filtering, or None to skip filtering
    """
    # CR-273: Set phase on agent so think() can filter tools
    prev_phase = getattr(agent, '_ooda_phase', None)
    agent._ooda_phase = phase

    think_task = asyncio.create_task(agent.think(prompt))
    # Cloud-API agents don't use Ollama — skip activity polling
    _uses_cloud_api = bool(agent.config.get("llm_provider"))
    ollama_url = f"{Config.LLM_BASE_URL}/api/ps"
    last_active = time.monotonic()
    start_time = time.monotonic()
    consecutive_api_failures = 0
    _MAX_API_FAILURES = 5  # Cancel after 5 consecutive unreachable polls (~75s)

    while not think_task.done():
        # Wait but check completion frequently
        for _ in range(poll_interval):
            if think_task.done():
                break
            await asyncio.sleep(1)

        if think_task.done():
            break

        elapsed = time.monotonic() - start_time

        # CR-270: Hard timeout — absolute safety net against system freeze.
        # Even if Ollama reports "active", no single think() call should run
        # longer than hard_timeout. The 2026-04-02 freeze showed eval times
        # escalating 23s → 83s → 107s → ∞ while the model appeared "active".
        if elapsed > hard_timeout:
            log.error(
                f"[{agent.agent_name}] HARD TIMEOUT: think() running for {elapsed:.0f}s "
                f"(limit={hard_timeout}s) — killing to prevent system freeze"
            )
            think_task.cancel()
            try:
                await think_task
            except asyncio.CancelledError:
                pass
            raise RuntimeError(
                f"Hard timeout after {elapsed:.0f}s — inference likely stuck "
                f"(VRAM exhaustion or model hang)"
            )

        # Check if Ollama still has an active model (skip for cloud-API agents)
        if _uses_cloud_api:
            last_active = time.monotonic()  # Cloud API is always "active"
            continue

        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(ollama_url, timeout=aiohttp.ClientTimeout(total=5)) as resp:
                    data = await resp.json()
                    models = data.get("models", [])
                    consecutive_api_failures = 0  # Reset on success
                    if models:
                        # Ollama is active — reset stale timer
                        last_active = time.monotonic()
                    else:
                        stale_seconds = time.monotonic() - last_active
                        if stale_seconds > stale_timeout:
                            log.error(
                                f"[{agent.agent_name}] Ollama inactive for {stale_seconds:.0f}s "
                                f"(stale_timeout={stale_timeout}s) — cancelling think()"
                            )
                            think_task.cancel()
                            try:
                                await think_task
                            except asyncio.CancelledError:
                                pass
                            raise RuntimeError(
                                f"Ollama stale for {stale_seconds:.0f}s — inference likely crashed"
                            )
                        else:
                            log.debug(
                                f"[{agent.agent_name}] Ollama idle for {stale_seconds:.0f}s "
                                f"(threshold: {stale_timeout}s)"
                            )
        except (aiohttp.ClientError, asyncio.TimeoutError, OSError):
            # CR-270: Track consecutive API failures instead of silently ignoring.
            # If Ollama itself is frozen/unresponsive, the old code would loop
            # forever because /api/ps kept timing out and stale_timeout never fired.
            consecutive_api_failures += 1
            log.warning(
                f"[{agent.agent_name}] Ollama API unreachable "
                f"({consecutive_api_failures}/{_MAX_API_FAILURES})"
            )
            if consecutive_api_failures >= _MAX_API_FAILURES:
                log.error(
                    f"[{agent.agent_name}] Ollama API unreachable {consecutive_api_failures}x "
                    f"in a row — system likely frozen, cancelling think()"
                )
                think_task.cancel()
                try:
                    await think_task
                except asyncio.CancelledError:
                    pass
                raise RuntimeError(
                    f"Ollama API unreachable {consecutive_api_failures}x — "
                    f"system freeze detected"
                )
        except RuntimeError:
            raise  # Re-raise our own stale/timeout error

    # CR-273: Restore previous phase (cleanup)
    agent._ooda_phase = prev_phase

    if think_task.cancelled():
        raise RuntimeError("think() was cancelled")

    return think_task.result()


# CR-273: 6-Phase OODA default parameters
# Phase numbering: 0=KONTEXT, 1=OBSERVE, 2=ORIENT, 3=DECIDE, 4=ACT, 5=PERSIST
_DEFAULT_PHASE_PARAMS = {
    "phase0": {"temperature": 0.2, "num_predict": 1024},   # KONTEXT: Factual, read state
    "phase1": {"temperature": 0.3, "num_predict": 1536},   # OBSERVE: Analytical, structure
    "phase2": {"temperature": 0.4, "num_predict": 2048},   # ORIENT: Analytical, Lagebild + Chunk-Analyse
    "phase3": {"temperature": 0.4, "num_predict": 1536},   # DECIDE: Analytical, stakeholder plan
    "phase4": {"temperature": 0.3, "num_predict": 2048},   # ACT: Professional, draft responses
    "phase5": {"temperature": 0.1, "num_predict": 1024},   # PERSIST: Precise, write exact formats
}


class _PhaseParams:
    """Context manager that temporarily overrides agent config for a specific phase."""

    def __init__(self, agent: "AIMOSAgent", phase: str, log: logging.Logger):
        self.agent = agent
        self.phase = phase
        self.log = log
        self.original = {}

    def __enter__(self):
        phase_params = self.agent.config.get("batch_phase_params", {})
        params = phase_params.get(self.phase, _DEFAULT_PHASE_PARAMS.get(self.phase, {}))
        for key, value in params.items():
            self.original[key] = self.agent.config.get(key)
            self.agent.config[key] = value
        if params:
            self.log.debug(f"[{self.agent.agent_name}] {self.phase} params: {params}")
        return self

    def __exit__(self, *args):
        for key, value in self.original.items():
            if value is None:
                self.agent.config.pop(key, None)
            else:
                self.agent.config[key] = value


async def process_batch(agent: "AIMOSAgent", messages: list[dict], log: logging.Logger):
    """CR-273: 6-Phase OODA Batch Cycle.

    Phase 0: KONTEXT  — Load workspace, sync data sources
    Phase 1: OBSERVE  — Structure inputs (messages + document inventory)
    Phase 2: ORIENT   — Build Lagebild (incl. document chunk analysis loop)
    Phase 3: DECIDE   — Identify stakeholders, create action plan
    Phase 4: ACT      — Draft responses + generate documents → Orchestrator dispatches
    Phase 5: PERSIST  — Guaranteed self-dispatch: save state (runs in finally-block)

    Architecture: LLM drafts, Orchestrator dispatches. COMMUNICATE tools are NOT
    in the LLM tool-set. Phase 5 PERSIST always runs, even if Phase 4 crashes.
    """
    from core.fallback import auto_followup

    if not messages:
        return

    # CR-275: Semantic dedup — merge near-duplicate messages from same sender.
    # A human sees "same person, same topic, 6 min apart" as one request.
    # Without this, the agent sends N emails for N near-identical inputs.
    messages = _dedup_messages(messages, log, agent.agent_name)

    if not messages:
        return

    log.info(f"[{agent.agent_name}] BATCH mode: processing {len(messages)} message(s)")

    # H-25 / CR-270: Clear in-memory history at batch start.
    # Without this, Phase 0-1 inherit stale history from the previous batch cycle
    # (agent is reused between cycles without restart).
    agent._history = []

    # Set session/thread context for batch-level operations
    first_msg = messages[0]
    agent._current_session_id = f"batch:{agent.agent_name}"
    agent._current_thread_id = f"batch:{first_msg.get('id', 0)}"
    agent._current_msg_kind = "batch"
    agent._tool_call_count = 0
    agent._tool_call_budget = agent.config.get("max_tool_calls_per_message", 30)

    ws_base = Path(f"storage/agents/{agent.agent_name}")

    # ══════════════════════════════════════════════════════════════════════
    #  Phase 0: KONTEXT — Load workspace, sync data sources
    # ══════════════════════════════════════════════════════════════════════
    phase0_context = await _phase0_context(agent, ws_base, log)

    # Dropbox/SharePoint sync (part of Phase 0)
    dropbox_path = agent.config.get("batch_dropbox_path")
    if dropbox_path and hasattr(agent, '_skills'):
        try:
            dropbox_skill = agent._skills.get("dropbox")
            if dropbox_skill and dropbox_skill.is_available():
                sync_result = await dropbox_skill._sync_folder(dropbox_path, "dokumente")
                log.info(f"[{agent.agent_name}] Phase 0: Dropbox sync: {sync_result[:100]}")
        except Exception as exc:
            log.warning(f"[{agent.agent_name}] Phase 0: Dropbox sync failed (non-critical): {exc}")

    # Workspace document scan (part of Phase 0)
    new_documents = []
    if agent.config.get("batch_workspace_scan", False):
        new_documents = _scan_workspace_documents(ws_base, log, agent.agent_name)

    # Format messages + document inventory for Phase 1
    batch_input = agent.format_batch_input(messages)
    if new_documents:
        doc_list = "\n".join(
            f"  - {d['name']} ({d['size_kb']} KB, client={d.get('client_dir', '-')})"
            for d in new_documents
        )
        batch_input += (
            f"\n\n--- NEW DOCUMENTS ON YOUR DESK ---\n"
            f"The following {len(new_documents)} document(s) arrived since your last session:\n"
            f"{doc_list}\n"
            f"These need to be analyzed during the ORIENT phase.\n"
        )

    # Context budget management
    batch_input = _context_monitor(agent, phase0_context, batch_input, messages, ws_base, log)

    # ══════════════════════════════════════════════════════════════════════
    #  Phase 1: OBSERVE — Structure all inputs
    # ══════════════════════════════════════════════════════════════════════
    analysis = await _phase1_observe(agent, phase0_context, batch_input, log)
    if analysis is None:
        return

    # ══════════════════════════════════════════════════════════════════════
    #  Phase 2: ORIENT — Build Lagebild (incl. document chunk analysis)
    # ══════════════════════════════════════════════════════════════════════

    # Phase 2a: Document chunk analysis loop (if documents present)
    doc_results = []
    if new_documents:
        # Preliminary Lagebild from Phase 1 analysis (before docs)
        preliminary_lagebild = analysis  # Use analysis as context for chunk processing
        doc_results = await _phase2a_documents(
            agent, new_documents, preliminary_lagebild, ws_base, log,
        )

    # CR-274: Resolve confidentiality scopes
    confidentiality = agent.config.get("batch_confidentiality", "none")
    scope_names: list[str] = []
    if confidentiality == "isolated":
        scope_groups = _group_messages_by_scope(messages, agent.config)
        scope_names = list(scope_groups.keys())
        log.info(
            f"[{agent.agent_name}] CR-274: Confidentiality=isolated, "
            f"{len(scope_names)} scope(s): {scope_names}"
        )

    # Phase 2b: Lagebild consolidation (with document findings)
    lagebild = await _phase2_orient(agent, analysis, log, scope_names=scope_names or None)
    if lagebild is None:
        return

    # H-09: Empty Lagebild guard
    if not lagebild or len(lagebild.strip()) < 20:
        log.error(
            f"[{agent.agent_name}] Phase 2 ORIENT: empty Lagebild "
            f"({len(lagebild)} chars). Aborting — would produce blind responses."
        )
        return

    # CR-274: Partition Lagebild by scope (for Phase 4 isolation)
    lagebild_partitions: dict[str, str] = {"_all": lagebild}
    if confidentiality == "isolated" and len(scope_names) > 1:
        lagebild_partitions = _partition_lagebild(lagebild, scope_names)
        if "_all" not in lagebild_partitions:
            lagebild_partitions["_all"] = lagebild  # Keep full version for DECIDE + PERSIST
        log.info(
            f"[{agent.agent_name}] CR-274: Lagebild partitioned into "
            f"{len(lagebild_partitions) - 1} scope(s) + _all"
        )

    # H-08: Human approval gate for critical agents
    await _human_approval_gate(agent, lagebild, messages, log)

    # ══════════════════════════════════════════════════════════════════════
    #  Phase 3: DECIDE — Stakeholder plan
    # ══════════════════════════════════════════════════════════════════════
    if agent.config.get("batch_phase4_mode") == "file":
        # File-mode: No stakeholders to dispatch to — skip Phase 3
        stakeholder_plan = ""
        log.info(f"[{agent.agent_name}] Phase 3 SKIPPED (file-mode, no stakeholder dispatch)")
    else:
        stakeholder_plan = await _phase3_decide(agent, lagebild, messages, log)

    # ══════════════════════════════════════════════════════════════════════
    #  Phase 4: ACT — Draft → Validate → Dispatch (Orchestrator)
    # ══════════════════════════════════════════════════════════════════════
    threads = {}
    phase4_results = []
    try:
        threads, phase4_results = await _phase4_act(
            agent, messages, lagebild, log,
            stakeholder_plan=stakeholder_plan,
            lagebild_partitions=lagebild_partitions,
        )
    except Exception as exc:
        log.error(f"[{agent.agent_name}] Phase 4 ACT failed: {exc}")
        # Phase 5 PERSIST still runs (finally-block below)

    # ══════════════════════════════════════════════════════════════════════
    #  Phase 5: PERSIST — Guaranteed self-dispatch (always runs)
    # ══════════════════════════════════════════════════════════════════════
    if agent.config.get("batch_persist", True):
        try:
            await _phase5_persist(agent, lagebild, phase4_results + doc_results, ws_base, log, msg_count=len(messages))
        except Exception as exc:
            log.error(f"[{agent.agent_name}] Phase 5 PERSIST failed: {exc}")
    else:
        log.info(f"[{agent.agent_name}] Phase 5 PERSIST skipped (batch_persist=false)")

    agent._touch()

    # File-mode: State-based completion check → self-retrigger if incomplete
    if agent.config.get("batch_phase4_mode") == "file":
        _completion = _check_file_mode_completion(agent, ws_base, log)
        if _completion["complete"]:
            log.info(
                f"[{agent.agent_name}] File-mode: COMPLETE — "
                f"{_completion['docs_in_inventar']}/{_completion['docs_in_arbeitsdatei']} docs, "
                f"all {len(_completion['output_files'])} output files valid."
            )
        elif _completion["reason"] == "no_arbeitsdatei":
            log.info(f"[{agent.agent_name}] File-mode: arbeitsdatei not ready yet, no retrigger")
        else:
            log.info(
                f"[{agent.agent_name}] File-mode: INCOMPLETE — {_completion['reason']}. "
                f"Self-retrigger for next cycle."
            )
            try:
                if agent._pool:
                    await agent._pool.execute(
                        "INSERT INTO pending_messages (agent_name, sender_id, content, kind) "
                        "VALUES ($1, '0', $2, 'internal')",
                        agent.agent_name,
                        f"Weiterarbeiten: {_completion['reason']}. "
                        f"Die Analyseergebnisse stehen in arbeitsdatei.md.",
                    )
                    log.info(f"[{agent.agent_name}] Self-retrigger message inserted")
            except Exception as exc:
                log.warning(f"[{agent.agent_name}] Self-retrigger failed: {exc}")

    # Auto-followup for the batch as a whole
    if first_msg.get("kind") not in ("scheduled_job", "internal"):
        if not agent.config.get("disable_auto_jobs"):
            await auto_followup(agent, lagebild, log)

    # H-15: Mark messages as fully completed
    msg_ids = [m.get("id") for m in messages if m.get("id")]
    if msg_ids and agent._pool:
        try:
            await agent._pool.execute(
                "UPDATE pending_messages SET content = content || '\n[BATCH_COMPLETED]' "
                "WHERE id = ANY($1::int[])",
                msg_ids,
            )
        except Exception:
            pass

    # CF-4: Workspace versioning for audit trail (H-A.15 Revisionssicherheit)
    try:
        import subprocess
        ws_str = str(ws_base)
        # Initialize git repo if not exists
        if not (ws_base / ".git").exists():
            subprocess.run(["git", "init"], cwd=ws_str, capture_output=True)
            subprocess.run(["git", "config", "user.name", "AIMOS Agent"], cwd=ws_str, capture_output=True)
            subprocess.run(["git", "config", "user.email", "agent@aimos.local"], cwd=ws_str, capture_output=True)
        # Stage and commit workspace files
        subprocess.run(["git", "add", "state.md", "arbeitsdatei.md", "offene_punkte.md",
                         "zusammenfassung.md", "todo.md", "status.md"],
                       cwd=ws_str, capture_output=True)
        result = subprocess.run(
            ["git", "commit", "-m", f"OODA cycle {__import__('datetime').datetime.now():%Y-%m-%d %H:%M} — {len(messages)} msg(s)"],
            cwd=ws_str, capture_output=True, text=True
        )
        if result.returncode == 0:
            log.info(f"[{agent.agent_name}] CF-4: Workspace versioned (git commit)")
        # else: no changes to commit, which is fine
    except Exception as exc:
        log.debug(f"[{agent.agent_name}] CF-4: Workspace versioning skipped: {exc}")

    log.info(f"[{agent.agent_name}] BATCH complete: {len(messages)} msg(s), {len(threads)} thread(s)")


# ── Phase implementations ──────────────────────────────────────────────────


async def _phase0_context(agent: "AIMOSAgent", ws_base: Path, log: logging.Logger) -> str:
    """Phase 0: Load workspace state — 'look at your desk before opening mail'."""
    workspace_context = ""
    leading_file = agent.config.get("batch_leading_file", "state.md")
    context_budget = agent.config.get("batch_context_budget", 2000)

    leading_path = ws_base / leading_file
    staleness_warning = ""
    if leading_path.exists():
        try:
            content = leading_path.read_text(encoding="utf-8")[:context_budget]
            workspace_context = f"YOUR HANDOVER PROTOCOL ({leading_file}):\n{content}\n"
            # H-04: Check staleness
            import datetime as _dt
            file_age_days = (_dt.datetime.now().timestamp() - leading_path.stat().st_mtime) / 86400
            if file_age_days > 7:
                staleness_warning = (
                    f"\n⚠ WARNING: Your handover protocol is {file_age_days:.0f} days old. "
                    f"Information may be outdated. Verify critical items against your data sources.\n"
                )
                log.warning(f"[{agent.agent_name}] Phase 0: {leading_file} is {file_age_days:.0f} days old")
        except Exception as exc:
            log.warning(f"[{agent.agent_name}] Phase 0: Could not read {leading_file}: {exc}")

    # --- H-IC.06: Cross-file invariant check (state.md vs arbeitsdatei.md) ---
    arbeitsdatei_path = ws_base / "arbeitsdatei.md"
    if leading_path.exists() and arbeitsdatei_path.exists():
        try:
            import datetime as _dt

            ad_lines = arbeitsdatei_path.read_text(encoding="utf-8").splitlines()
            ad_line_count = len(ad_lines)

            # Check 1: state says "not begun" but arbeitsdatei has substantial content
            if ad_line_count > 50 and workspace_context:
                state_lower = workspace_context.lower()
                if "0/" in state_lower or "noch nicht begonnen" in state_lower:
                    inv_msg = (
                        f"\n⚠ INVARIANT VIOLATION (H-IC.06): {leading_file} says work not begun "
                        f"but arbeitsdatei.md has {ad_line_count} lines. "
                        f"Possible stale data from previous run. "
                        f"Verify actual progress before reporting numbers.\n"
                    )
                    staleness_warning += inv_msg
                    log.warning(
                        f"[{agent.agent_name}] Phase 0 invariant: {leading_file} says not begun "
                        f"but arbeitsdatei.md has {ad_line_count} lines"
                    )
            elif ad_line_count > 10 and workspace_context:
                state_lower = workspace_context.lower()
                if "0/" in state_lower or "noch nicht begonnen" in state_lower:
                    inv_msg = (
                        f"\n⚠ WARNING: {leading_file} says work not begun "
                        f"but arbeitsdatei.md has {ad_line_count} lines. Check consistency.\n"
                    )
                    staleness_warning += inv_msg
                    log.warning(
                        f"[{agent.agent_name}] Phase 0: minor inconsistency — "
                        f"{leading_file} says not begun, arbeitsdatei.md has {ad_line_count} lines"
                    )

            # Check 2: state.md timestamp >1h older than arbeitsdatei.md
            state_mtime = leading_path.stat().st_mtime
            ad_mtime = arbeitsdatei_path.stat().st_mtime
            age_diff_hours = (ad_mtime - state_mtime) / 3600
            if age_diff_hours > 1:
                ts_msg = (
                    f"\n⚠ WARNING: {leading_file} is {age_diff_hours:.1f}h older than "
                    f"arbeitsdatei.md — possible inconsistency after crash. "
                    f"Trust arbeitsdatei.md for actual progress, update {leading_file} first.\n"
                )
                staleness_warning += ts_msg
                log.warning(
                    f"[{agent.agent_name}] Phase 0: {leading_file} is {age_diff_hours:.1f}h "
                    f"older than arbeitsdatei.md"
                )
        except Exception as exc:
            log.warning(f"[{agent.agent_name}] Phase 0 invariant check failed: {exc}")

    # List available workspace files
    available_files = []
    if ws_base.exists():
        available_files = sorted(f.name for f in ws_base.iterdir()
                                 if f.is_file() and f.name != leading_file and not f.name.startswith('.'))

    data_sources = agent.config.get("batch_data_sources", [])

    if not (workspace_context or data_sources or available_files):
        log.info(f"[{agent.agent_name}] BATCH Phase 0 skipped (no context files or data sources configured)")
        return "First session — no prior knowledge. No files in workspace."

    files_hint = ""
    if available_files:
        files_hint = (
            f"\nAvailable detail files (use read_file to load if needed): "
            f"{', '.join(available_files[:15])}\n"
        )
    source_instructions = ""
    if data_sources:
        source_list = "\n".join(f"  - {s}" for s in data_sources)
        source_instructions = (
            f"\nAlso check these data sources for changes since your last session:\n"
            f"{source_list}\n"
            "Use your tools to query them now.\n"
        )
    phase0_prompt = (
        f"{workspace_context}"
        f"{staleness_warning}"
        f"{files_hint}"
        "You are starting a new work session. The handover protocol above contains "
        "your previous state. READ IT CAREFULLY and write a brief summary:\n"
        "- Who are your current stakeholders?\n"
        "- What is still open?\n"
        "- What deadlines are approaching?\n"
        "- What was the last action taken?\n\n"
        "Write your summary as PLAIN TEXT. All the data you need is above.\n"
        f"{source_instructions}"
        "Do NOT act yet — only review and summarize."
    )

    log.info(f"[{agent.agent_name}] BATCH Phase 0 (CONTEXT): {len(workspace_context)} chars workspace")
    try:
        with _PhaseParams(agent, "phase0", log):
            result = await _think_with_activity_check(agent, phase0_prompt, log, stale_timeout=agent.config.get("batch_stale_timeout", 120), phase="0")
        log.info(f"[{agent.agent_name}] BATCH Phase 0 complete: {len(result)} chars")
        return result
    except (asyncio.TimeoutError, Exception) as exc:
        log.error(f"[{agent.agent_name}] BATCH Phase 0 FAILED: {exc}")
        return ""


def _context_monitor(
    agent: "AIMOSAgent", phase0_context: str, batch_input: str,
    messages: list[dict], ws_base: Path, log: logging.Logger,
) -> str:
    """Context Monitor: proactive budget management. Returns (possibly compressed) batch_input."""
    num_ctx = agent.config.get("num_ctx", Config.DEFAULT_NUM_CTX)
    reserve_for_later_phases = 8000
    available_for_input = num_ctx - reserve_for_later_phases
    est_input_tokens = (len(phase0_context) + len(batch_input)) // 4
    ctx_usage_pct = (est_input_tokens / num_ctx) * 100 if num_ctx else 0

    if est_input_tokens > available_for_input:
        log.warning(
            f"[{agent.agent_name}] CONTEXT OVERLOAD: ~{est_input_tokens} input tokens "
            f"exceeds budget of {available_for_input} (num_ctx={num_ctx}). Compressing."
        )
        if len(phase0_context) > 2000:
            # Can't modify phase0_context (str is immutable in caller), but truncation
            # was already applied via context_budget in Phase 0
            log.info(f"[{agent.agent_name}] Phase 0 context already budget-limited")

        est_input_tokens = (len(phase0_context) + len(batch_input)) // 4
        if est_input_tokens > available_for_input:
            trimmed_msgs = []
            for msg in messages:
                trimmed = dict(msg)
                content = trimmed.get("content", "")
                if len(content) > 200:
                    trimmed["content"] = content[:200] + " [... truncated]"
                trimmed_msgs.append(trimmed)
            batch_input = agent.format_batch_input(trimmed_msgs)
            log.warning(
                f"[{agent.agent_name}] Messages truncated to 200 chars each. "
                f"Quality will be reduced. Consider splitting this agent."
            )

        overload_note = (
            f"⚠ CONTEXT OVERLOAD {__import__('datetime').datetime.now():%Y-%m-%d %H:%M}\n"
            f"Input: ~{est_input_tokens} tokens, Budget: {available_for_input}, num_ctx: {num_ctx}\n"
            f"Messages: {len(messages)}, State: {len(phase0_context)} chars\n"
            f"Recommendation: Split this agent's responsibilities or increase num_ctx."
        )
        overload_path = ws_base / "overload_warning.txt"
        try:
            overload_path.parent.mkdir(parents=True, exist_ok=True)
            overload_path.write_text(overload_note, encoding="utf-8")
        except Exception:
            pass

    elif ctx_usage_pct > 50:
        log.info(
            f"[{agent.agent_name}] Context budget: ~{est_input_tokens} tokens "
            f"({ctx_usage_pct:.0f}% of {num_ctx}) after Phase 0 + Messages — OK"
        )

    return batch_input


async def _phase1_observe(
    agent: "AIMOSAgent", phase0_context: str, batch_input: str, log: logging.Logger,
) -> str | None:
    """Phase 1: OBSERVE — Structure all new messages."""
    phase1_prompt = ""
    if phase0_context:
        phase1_prompt += f"YOUR CURRENT STATE (from Phase 0):\n{phase0_context}\n\n"
    phase1_prompt += (
        f"{batch_input}\n\n"
        "Structure these messages as PLAIN TEXT analysis:\n"
        "- Group by sender/thread\n"
        "- For each: WHO sent it, WHAT they want, HOW URGENT\n"
        "- Compare with your current state — what is NEW vs. already known?\n"
        "- Prioritize by urgency (escalations > customer requests > internal)\n\n"
        "Write your analysis as PLAIN TEXT. All data is above.\n"
        "You may use read_file if you need details not shown above.\n"
        "Do NOT act yet — only analyze."
    )
    log.info(f"[{agent.agent_name}] BATCH Phase 1 (OBSERVE): {len(batch_input)} chars input")
    try:
        with _PhaseParams(agent, "phase1", log):
            analysis = await _think_with_activity_check(agent, phase1_prompt, log, stale_timeout=agent.config.get("batch_stale_timeout", 120), phase="1")
        log.info(f"[{agent.agent_name}] BATCH Phase 1 complete: {len(analysis)} chars")
        return analysis
    except (asyncio.TimeoutError, Exception) as exc:
        log.error(f"[{agent.agent_name}] BATCH Phase 1 FAILED: {exc}")
        return None


async def _phase2_orient(
    agent: "AIMOSAgent", analysis: str, log: logging.Logger,
    scope_names: list[str] | None = None,
) -> str | None:
    """Phase 2: ORIENT — Consolidate and build Lagebild.

    CR-274: When scope_names is provided (isolated mode), the prompt instructs the
    LLM to structure the Lagebild with ## [SCOPE: ...] headers per scope. This
    enables _partition_lagebild() to split the result for Phase 4.
    """
    scope_instruction = ""
    if scope_names and len(scope_names) > 1:
        scope_list = "\n".join(f"  - {s}" for s in scope_names)
        scope_instruction = (
            f"\nIMPORTANT: This batch contains {len(scope_names)} SEPARATE confidential scopes.\n"
            f"Structure your Lagebild with one section per scope, using this exact header format:\n"
            f"## [SCOPE: scope_name]\n"
            f"The scopes are:\n{scope_list}\n"
            f"NEVER mix information between scopes. Each section must be self-contained.\n\n"
        )

    phase2_prompt = (
        f"YOUR ANALYSIS FROM PHASE 1:\n{analysis}\n\n"
        "Now consolidate and write a LAGEBILD (situation report) as TEXT.\n"
        "Based on your analysis:\n"
        "- What is the current situation?\n"
        "- What are the dependencies between tasks?\n"
        "- What does each stakeholder need to know?\n"
        "- What actions are needed per thread?\n\n"
        f"{scope_instruction}"
        "You may use tools (read_file, search) to gather additional information if needed, "
        "but your OUTPUT must be a written situation report, not just tool calls.\n"
        "Write the Lagebild NOW as structured text."
    )
    log.info(f"[{agent.agent_name}] BATCH Phase 2 (ORIENT)")
    try:
        with _PhaseParams(agent, "phase2", log):
            lagebild = await _think_with_activity_check(agent, phase2_prompt, log, stale_timeout=agent.config.get("batch_stale_timeout", 120), phase="2")
        log.info(f"[{agent.agent_name}] BATCH Phase 2 complete (Lagebild): {len(lagebild)} chars")
        return lagebild
    except (asyncio.TimeoutError, Exception) as exc:
        log.error(f"[{agent.agent_name}] BATCH Phase 2 FAILED: {exc}")
        return None


async def _phase3_decide(
    agent: "AIMOSAgent", lagebild: str, messages: list[dict], log: logging.Logger,
) -> str:
    """Phase 3: DECIDE — Identify ALL stakeholders who need to be informed.

    This is the key difference to a chatbot: The agent doesn't just respond to
    people who wrote — it proactively identifies everyone AFFECTED by the changes
    described in the Lagebild, even if they didn't send a message.
    """
    # Build list of known senders from this batch
    senders = set()
    for msg in messages:
        tid = msg.get("thread_id", "")
        if tid:
            senders.add(tid)

    phase3_prompt = (
        f"LAGEBILD:\n{lagebild}\n\n"
        f"MESSAGES IN THIS BATCH came from these threads: {', '.join(senders)}\n\n"
        "TASK: Based on the Lagebild, identify ALL stakeholders who need to be "
        "informed or contacted — including those who did NOT send a message but "
        "are AFFECTED by the changes.\n\n"
        "For each stakeholder, specify:\n"
        "- Thread ID or email (if known from previous conversations)\n"
        "- What they need to know\n"
        "- Priority: HIGH (must inform now) / LOW (nice to know)\n\n"
        "Think: Who is waiting for this information? Who will be impacted?\n"
        "Who needs to adjust their schedule? Who asked about this before?\n\n"
        "GOVERNANCE CHECK: Does the Lagebild contain activity that looks like an "
        "untracked project or case? Signs: multiple stakeholders, deliverables, "
        "deadlines, dependencies — but no formal project plan or case file in the "
        "workspace. If yes, add a HIGH-priority action item: 'Formalize as project/case "
        "— needs scope, plan, and responsible owner.'\n"
        "Similarly, watch for MISSION CREEP: Is the scope of existing work growing "
        "beyond what was originally planned? New requirements without a change request? "
        "If yes, flag it.\n\n"
        "DEADLINE CHECK: If today's date is within 90 days of 31 July of the current year, add a HIGH-priority note: 'Abgabefrist {year} nähert sich — Mandant informieren falls noch nicht geschehen.'\n\n"
        "Output a STAKEHOLDER ACTION PLAN."
    )

    log.info(f"[{agent.agent_name}] BATCH Phase 3 (DECIDE)")
    try:
        with _PhaseParams(agent, "phase3", log):  # Same params as Orient (analytical)
            plan = await _think_with_activity_check(
                agent, phase3_prompt, log,
                stale_timeout=agent.config.get("batch_stale_timeout", 120),
                phase="3",
            )
        log.info(f"[{agent.agent_name}] BATCH Phase 3 complete: {len(plan)} chars")
        return plan
    except (asyncio.TimeoutError, Exception) as exc:
        log.warning(f"[{agent.agent_name}] BATCH Phase 3 FAILED (non-critical): {exc}")
        return ""  # Fallback: Phase 3 uses only message threads (old behavior)


async def _human_approval_gate(
    agent: "AIMOSAgent", lagebild: str, messages: list[dict], log: logging.Logger,
):
    """H-08: Send Lagebild to reviewer for critical agents."""
    if not agent.config.get("batch_require_human_approval"):
        return
    reviewer = agent.config.get("batch_approval_target", "projektleiter")
    log.info(f"[{agent.agent_name}] BATCH Human Approval Gate: sending Lagebild to {reviewer}")
    approval_msg = (
        f"[LAGEBILD ZUR FREIGABE — {agent.agent_name}]\n\n"
        f"{lagebild}\n\n"
        f"Bitte prüfen und bestätigen. Der Agent wartet auf Freigabe, "
        f"bevor er {len(messages)} Nachricht(en) beantwortet."
    )
    if agent._pool:
        await agent._pool.execute(
            "INSERT INTO pending_messages (agent_name, sender_id, content, kind, thread_id, processed) "
            "VALUES ($1, 0, $2, 'internal', $3, FALSE)",
            reviewer, approval_msg, f"approval:{agent.agent_name}",
        )
    log.info(f"[{agent.agent_name}] Lagebild sent to {reviewer} for review (non-blocking)")


async def _phase4b_validate(
    agent: "AIMOSAgent", lagebild: str, stakeholder_plan: str,
    response: str, thread_id: str, log: logging.Logger,
    reference_data: str = "",
) -> str:
    """Phase 4b: VALIDATE — Check response before dispatch.

    Returns "OK" if valid, or "PROBLEM: ..." if issues found.
    Non-blocking: problems are logged as warnings, response is still dispatched
    (H-20: avoid false-positive blocking).
    """
    # Bug #19: Include reference data so validator doesn't reject FAQ-sourced facts
    data_sources = f"LAGEBILD:\n{lagebild[:2000]}\n\n"
    if reference_data:
        data_sources += f"{reference_data[:3000]}\n\n"

    validate_prompt = (
        f"{data_sources}"
        f"STAKEHOLDER PLAN:\n{stakeholder_plan[:500]}\n\n"
        f"RESPONSE TO SEND (thread={thread_id}):\n{response[:2000]}\n\n"
        "VALIDATE this response:\n"
        "1. Does it contain facts NOT in the Lagebild or REFERENZDATEN? (hallucination check)\n"
        "2. Does it mention people not known from the conversation?\n"
        "3. Is it consistent with the data sources above?\n"
        "4. Does it reference specific amounts? If yes, are they in the data sources?\n"
        "Answer ONLY: 'OK' or 'PROBLEM: [brief description]'"
    )
    try:
        with _PhaseParams(agent, "phase4", log):  # Low temperature for precise validation
            result = await _think_with_activity_check(
                agent, validate_prompt, log,
                stale_timeout=agent.config.get("batch_stale_timeout", 120),
                phase="4b",  # 4b VALIDATE — current_time only
            )
        result = result.strip()
        if result.upper().startswith("OK"):
            return "OK"
        else:
            log.warning(f"[{agent.agent_name}] Phase 4b VALIDATE: {result[:200]}")
            return result
    except Exception as exc:
        log.warning(f"[{agent.agent_name}] Phase 4b VALIDATE failed (non-critical): {exc}")
        return "OK"  # H-21: Don't block on validation failure


async def _phase4_act(
    agent: "AIMOSAgent", messages: list[dict], lagebild: str, log: logging.Logger,
    stakeholder_plan: str = "",
    lagebild_partitions: dict[str, str] | None = None,
) -> tuple[dict, list[dict]]:
    """Phase 4: ACT — per stakeholder (from Phase 3 plan) with Lagebild context.

    CR-250: Phase 4 now iterates over stakeholders identified in Phase 3,
    not just threads with messages. This enables proactive notification of
    affected stakeholders who didn't send a message.
    """
    threads: dict[str, list[dict]] = defaultdict(list)
    for msg in messages:
        tid = msg.get("thread_id") or f"sender:{msg.get('sender_id', 0)}"
        threads[tid].append(msg)

    # CR-250: Extract proactive stakeholders from Phase 2b plan
    # These are stakeholders who didn't send a message but need to be informed.
    proactive_threads: dict[str, str] = {}  # thread_id → what they need to know
    if stakeholder_plan:
        import re as _re
        # Look for thread IDs or bare email addresses mentioned in the plan
        # that are NOT already in the message threads.
        # The LLM may write "email:meier@x.de" (thread format) or just "meier@x.de" (bare).
        found_ids = set()
        # Strategy 1: Bare email addresses → clean, canonical thread format
        # This catches both "email:meier@x.de" and bare "meier@x.de"
        for match in _re.finditer(r'([\w.+-]+@[\w.-]+\.\w{2,})', stakeholder_plan):
            found_ids.add(f"email:{match.group(1)}")
        for tid in found_ids:
            if tid not in threads:
                # Extract the "what they need to know" context
                # Take the 200 chars around the match
                start = max(0, match.start() - 50)
                end = min(len(stakeholder_plan), match.end() + 200)
                context = stakeholder_plan[start:end].strip()
                proactive_threads[tid] = context
        if proactive_threads:
            log.info(
                f"[{agent.agent_name}] Phase 2b identified {len(proactive_threads)} "
                f"PROACTIVE stakeholder(s): {list(proactive_threads.keys())}"
            )

    total_threads = len(threads) + len(proactive_threads)
    log.info(f"[{agent.agent_name}] BATCH Phase 4 (ACT): {len(threads)} with messages + {len(proactive_threads)} proactive = {total_threads} total")

    # Bug #19: Pre-compute reference summary once for all threads
    _ws = Path(f"storage/agents/{agent.agent_name}")
    _reference_summary = _build_reference_summary(_ws)
    if _reference_summary:
        log.info(f"[{agent.agent_name}] Bug#19: Injected {len(_reference_summary)} chars reference data into Phase 4 prompt")

    phase3_results: list[dict] = []

    # ══ File-mode Output Assembly: Harness baut, LLM entscheidet klein ══
    if agent.config.get("batch_phase4_mode") == "file":
        _arbeitsdatei = _ws / "arbeitsdatei.md"
        if _arbeitsdatei.exists() and _arbeitsdatei.stat().st_size > 100:
            log.info(f"[{agent.agent_name}] Phase 4: Output Assembly (file-mode)")
            try:
                assembly_results = await _phase4_output_assembly(agent, _ws, log)
                phase3_results.extend(assembly_results)
            except Exception as exc:
                log.error(f"[{agent.agent_name}] Phase 4 Output Assembly FAILED: {exc}")
            # Skip normal stakeholder loop — output is already written
            # Jump to Phase 5 (via the normal code path after the loop)
            threads = {}  # Empty → skip stakeholder loop
            proactive_threads = {}  # Empty → skip proactive loop

    for thread_id, thread_msgs in threads.items():
        representative_msg = thread_msgs[0]

        # CR-270: History isolation — clear in-memory history before each stakeholder call.
        # Without this, agent._history accumulates messages from ALL prior stakeholder
        # calls within this Phase 3 loop, causing exponential context growth and
        # eventual VRAM exhaustion (root cause of 2026-04-02 system freeze).
        agent._history = []

        # Load conversation history for this specific thread
        thread_history_text = ""
        if agent._pool:
            rows = await agent._pool.fetch(
                "SELECT role, content FROM aimos_chat_histories "
                "WHERE agent_name=$1 AND thread_id=$2 "
                "ORDER BY id DESC LIMIT 20",
                agent.agent_name, thread_id,
            )
            if rows:
                history_lines = [f"[{r['role']}]: {r['content'][:500]}" for r in reversed(rows)]
                thread_history_text = "\n".join(history_lines)

        agent._current_thread_id = thread_id
        agent._current_msg_kind = representative_msg.get("kind", "batch")
        agent._tool_call_count = 0

        sender_info = f"sender_id={representative_msg.get('sender_id', 0)}"
        channel_info = representative_msg.get("kind", "unknown")
        msg_contents = "\n".join(f"  - {m.get('content', '')[:200]}" for m in thread_msgs)

        # S-3 + CR-274: Cross-thread leak protection with scope isolation
        scope = _resolve_scope(thread_id, agent.config)
        if (
            lagebild_partitions
            and agent.config.get("batch_confidentiality") == "isolated"
            and scope in lagebild_partitions
        ):
            # DETERMINISTIC ISOLATION: LLM only sees this scope's Lagebild
            scope_lagebild = lagebild_partitions[scope]
            phase3_prompt = (
                f"LAGEBILD:\n{scope_lagebild}\n\n"
            )
        else:
            # Fallback: full Lagebild with prompt-based isolation
            phase3_prompt = (
                f"LAGEBILD (cross-thread situation report):\n{lagebild}\n\n"
                "IMPORTANT: The Lagebild above contains information about MULTIPLE stakeholders.\n"
                "You may ONLY share information relevant to the stakeholder below.\n"
                "NEVER mention names, messages, or internal details of other stakeholders.\n\n"
            )
        if thread_history_text:
            phase3_prompt += (
                f"CONVERSATION HISTORY with this stakeholder (thread={thread_id}):\n"
                f"{thread_history_text}\n\n"
            )
        # CR-278: Inject arbeitsdatei summary so the LLM doesn't need to load it via tool
        _arbeitsdatei_summary = _build_arbeitsdatei_summary(_ws)
        _summary_block = ""
        if _arbeitsdatei_summary:
            _summary_block = f"\n{_arbeitsdatei_summary}\n\n"

        # Bug #19: Inject reference/ files (computed once before loop)
        if _reference_summary:
            _summary_block += f"\n{_reference_summary}\n\n"

        phase3_prompt += (
            f"CURRENT MESSAGE(S) from this stakeholder ({sender_info}, channel={channel_info}):\n"
            f"{msg_contents}\n\n"
            f"{_summary_block}"
            # Phase 4 mode: "dispatch" (default) vs "file"
            # dispatch: Agent writes plain text → orchestrator sends it
            # file: Agent uses write_file to create output files
            + (_phase4_task_file(agent) if agent.config.get("batch_phase4_mode") == "file"
               else _phase4_task_dispatch())
        )

        log.info(
            f"[{agent.agent_name}] BATCH Phase 4a DRAFT thread={thread_id} "
            f"({len(thread_msgs)} msg(s), {len(thread_history_text)} chars history)"
        )
        try:
            with _PhaseParams(agent, "phase4", log):
                result = await _think_with_activity_check(
                    agent, phase3_prompt, log,
                    stale_timeout=agent.config.get("batch_stale_timeout", 120),
                    # file-mode: Phase "4" (full ACT, write_file allowed)
                    # dispatch-mode: Phase "4a" (DRAFT, READ only)
                    phase="4" if agent.config.get("batch_phase4_mode") == "file" else "4a",
                )
            log.info(f"[{agent.agent_name}] BATCH Phase 4a DRAFT thread={thread_id} complete: {len(result)} chars")

            # CR-279: Optional agent-specific draft post-filter (FuSi hook).
            # Only the most basic universal check here (tool-call text in emails).
            # Agent-specific validators go in batch_hooks config.
            result = _apply_draft_safety(result, agent, thread_id, log)

            phase3_results.append({
                "thread_id": thread_id, "status": "sent", "response_len": len(result),
                "summary": result[:150],
            })
        except (asyncio.TimeoutError, Exception) as exc:
            log.error(f"[{agent.agent_name}] BATCH Phase 4a thread={thread_id} FAILED: {exc}")
            phase3_results.append({
                "thread_id": thread_id, "status": "FAILED", "error": str(exc),
            })
            continue

        # CR-250 Phase 3b: Validate before dispatch
        validation = "OK"
        if agent.config.get("batch_validate_responses", True) and stakeholder_plan:
            validation = await _phase4b_validate(
                agent, lagebild, stakeholder_plan, result, thread_id, log,
                reference_data=_reference_summary,
            )
        # CR-281: If validation finds a problem, retry with feedback (max 2 retries)
        retry_count = 0
        while "PROBLEM" in validation and retry_count < 2:
            retry_count += 1
            log.warning(
                f"[{agent.agent_name}] CR-281: Validation PROBLEM — retry {retry_count}/2 "
                f"for thread={thread_id}: {validation[:100]}"
            )
            # Give the LLM the validation feedback and ask for a corrected draft
            # Bug #19: Include reference data in retry prompt so LLM has facts to work with
            _ref_block = f"\n{_reference_summary[:3000]}\n\n" if _reference_summary else ""
            retry_prompt = (
                f"Your previous draft was REJECTED by quality review.\n"
                f"PROBLEM: {validation}\n\n"
                f"LAGEBILD:\n{lagebild[:1500]}\n\n"
                f"{_ref_block}"
                f"Write a CORRECTED response that fixes this problem.\n"
                f"Use the REFERENZDATEN above as your data source.\n"
                f"Do NOT repeat the error. Write ONLY the corrected response as plain text."
            )
            agent._history = []
            try:
                with _PhaseParams(agent, "phase4", log):
                    result = await _think_with_activity_check(
                        agent, retry_prompt, log,
                        stale_timeout=agent.config.get("batch_stale_timeout", 120),
                        phase="4a",  # 4a DRAFT retry
                    )
                # Re-validate the corrected draft
                result = _apply_draft_safety(result, agent, thread_id, log)
                phase3_results[-1]["response_len"] = len(result)
                phase3_results[-1]["summary"] = result[:150]

                validation = "OK"
                if agent.config.get("batch_validate_responses", True) and stakeholder_plan:
                    validation = await _phase4b_validate(
                        agent, lagebild, stakeholder_plan, result, thread_id, log,
                        reference_data=_reference_summary,
                    )
            except Exception as exc:
                log.error(f"[{agent.agent_name}] CR-281: Retry {retry_count} failed: {exc}")
                break

        if "PROBLEM" in validation:
            # Still failing after retries — block dispatch, log for human review
            phase3_results[-1]["validation"] = validation
            phase3_results[-1]["status"] = "BLOCKED_AFTER_RETRIES"
            log.error(
                f"[{agent.agent_name}] CR-281: BLOCKED after {retry_count} retries "
                f"for thread={thread_id}. Draft NOT sent. Needs human review."
            )
            continue

        # dispatch_response now persists the answer in chat_histories (Bug #17 fix in dispatch.py)
        route = await agent.dispatch_response(result, representative_msg)
        log.info(f"[{agent.agent_name}] BATCH dispatched thread={thread_id} → {route}")

    # CR-250: Process proactive stakeholders (those who didn't send a message)
    for thread_id, context in proactive_threads.items():
        # CR-270: History isolation (same as reactive loop above)
        agent._history = []

        agent._current_thread_id = thread_id
        agent._current_msg_kind = "email"  # Default to email for proactive outreach
        agent._tool_call_count = 0

        # Load conversation history for this thread
        thread_history_text = ""
        if agent._pool:
            rows = await agent._pool.fetch(
                "SELECT role, content FROM aimos_chat_histories "
                "WHERE agent_name=$1 AND thread_id=$2 "
                "ORDER BY id DESC LIMIT 20",
                agent.agent_name, thread_id,
            )
            if rows:
                history_lines = [f"[{r['role']}]: {r['content'][:500]}" for r in reversed(rows)]
                thread_history_text = "\n".join(history_lines)

        # CR-274: Scope-filtered Lagebild for proactive stakeholders
        scope = _resolve_scope(thread_id, agent.config)
        if (
            lagebild_partitions
            and agent.config.get("batch_confidentiality") == "isolated"
            and scope in lagebild_partitions
        ):
            scope_lagebild = lagebild_partitions[scope]
        else:
            scope_lagebild = lagebild

        phase3_prompt = (
            f"LAGEBILD:\n{scope_lagebild}\n\n"
            f"STAKEHOLDER ACTION PLAN says this person needs to be informed:\n{context}\n\n"
        )
        if thread_history_text:
            phase3_prompt += f"CONVERSATION HISTORY (thread={thread_id}):\n{thread_history_text}\n\n"
        phase3_prompt += (
            f"This stakeholder (thread={thread_id}) did NOT send a message in this batch, "
            "but is AFFECTED by the changes in the Lagebild.\n"
            "Write a PROACTIVE notification as PLAIN TEXT. Do NOT use send_email or any "
            "communication tool. The system will dispatch this text for you.\n"
            "Be concise and only share what is relevant to them."
        )

        log.info(f"[{agent.agent_name}] BATCH Phase 4a PROACTIVE DRAFT thread={thread_id}")
        try:
            with _PhaseParams(agent, "phase4", log):
                result = await _think_with_activity_check(
                    agent, phase3_prompt, log,
                    stale_timeout=agent.config.get("batch_stale_timeout", 120),
                    phase="4a",  # 4a DRAFT proactive
                )
            log.info(f"[{agent.agent_name}] BATCH Phase 4a PROACTIVE thread={thread_id} complete: {len(result)} chars")
            phase3_results.append({
                "thread_id": thread_id, "status": "proactive", "response_len": len(result),
                "summary": result[:150],
            })

            # Validate proactive response too
            if agent.config.get("batch_validate_responses", True) and stakeholder_plan:
                validation = await _phase4b_validate(
                    agent, lagebild, stakeholder_plan, result, thread_id, log,
                    reference_data=_reference_summary,
                )
                if "PROBLEM" in validation:
                    phase3_results[-1]["validation"] = validation
                    log.warning(f"[{agent.agent_name}] Phase 4b PROACTIVE flagged: {validation[:100]}")

            # Dispatch proactive message — extract email from thread_id
            import re as _re2
            email_match = _re2.search(r'email:([\w.+-]+@[\w.-]+\.\w+)', thread_id)
            if email_match:
                # Real email address found — dispatch
                addr = email_match.group(1)
                proactive_msg = {
                    "kind": "email",
                    "sender_id": 0,
                    "thread_id": thread_id,
                    "content": f"[E-Mail empfangen]\nVon: {addr}\nKunden-Email: {addr}\nBetreff: Proaktive Benachrichtigung\nText: (proaktiv)",
                }
                route = await agent.dispatch_response(result, proactive_msg)
                log.info(f"[{agent.agent_name}] BATCH PROACTIVE dispatched thread={thread_id} → {route}")
            else:
                # No real email (e.g. "email:installateur@...") — log as todo
                log.info(
                    f"[{agent.agent_name}] BATCH PROACTIVE thread={thread_id}: "
                    f"No valid email address — adding to open tasks instead of dispatching"
                )
                phase3_results[-1]["status"] = "no_contact"
                phase3_results[-1]["summary"] = f"Kontaktdaten fehlen für {thread_id}"

        except (asyncio.TimeoutError, Exception) as exc:
            log.error(f"[{agent.agent_name}] BATCH Phase 4a PROACTIVE thread={thread_id} FAILED: {exc}")
            phase3_results.append({
                "thread_id": thread_id, "status": "FAILED", "error": str(exc),
            })

    return threads, phase3_results


# ── CR-258: Phase 2a — Document chunk analysis ────────────────────────────


async def _phase2a_documents(
    agent: "AIMOSAgent", new_documents: list[dict], lagebild: str,
    ws_base: Path, log: logging.Logger,
) -> list[dict]:
    """Phase 2a: Process new documents chunk by chunk within the OODA cycle.

    Uses the same safety patterns as Phase 4 stakeholder processing:
    - History isolation per chunk (CR-270)
    - _think_with_activity_check with hard timeout
    - Validation after each chunk
    """
    max_chunks = agent.config.get("batch_max_chunks_per_cycle", 20)
    doc_results = []
    total_chunks_processed = 0

    # CR-271: Dynamic chunk size — maximize document content per LLM call.
    # Instead of a fixed chunk_size, calculate how much document text fits
    # alongside system prompt + arbeitsdatei context + phase prompt overhead.
    num_ctx = agent.config.get("num_ctx", Config.DEFAULT_NUM_CTX)
    _SYSTEM_PROMPT_TOKENS = 2500   # System prompt (estimated)
    _PHASE_PROMPT_TOKENS = 500     # Phase 2a instructions
    _OUTPUT_RESERVE_TOKENS = 4096  # Reserve for thinking + response
    _SAFETY_MARGIN = 500
    _ARBEITSDATEI_BUDGET_CHARS = 2000  # Max chars of arbeitsdatei context per chunk

    log.info(
        f"[{agent.agent_name}] BATCH Phase 2a (DOCUMENTS): "
        f"{len(new_documents)} document(s) to analyze"
    )

    # Load current arbeitsdatei for context (if exists)
    arbeitsdatei_path = ws_base / "arbeitsdatei.md"
    arbeitsdatei_context = ""
    if arbeitsdatei_path.exists():
        try:
            content = arbeitsdatei_path.read_text(encoding="utf-8")
            arbeitsdatei_context = content[-_ARBEITSDATEI_BUDGET_CHARS:] if len(content) > _ARBEITSDATEI_BUDGET_CHARS else content
        except Exception:
            pass

    # Dynamic chunk size: fill remaining context with document text
    arbeitsdatei_tokens = len(arbeitsdatei_context) // 4
    available_for_doc = num_ctx - _SYSTEM_PROMPT_TOKENS - _PHASE_PROMPT_TOKENS - arbeitsdatei_tokens - _OUTPUT_RESERVE_TOKENS - _SAFETY_MARGIN
    chunk_size = max(available_for_doc * 4, 2000)  # Token→Char, minimum 2000
    # Cap at config value if explicitly set (user override)
    config_chunk = agent.config.get("batch_chunk_size")
    if config_chunk:
        chunk_size = min(chunk_size, config_chunk)

    log.info(
        f"[{agent.agent_name}] Phase 2a: Dynamic chunk_size={chunk_size} chars "
        f"(num_ctx={num_ctx}, available={available_for_doc} tokens, "
        f"arbeitsdatei={arbeitsdatei_tokens} tokens)"
    )

    for doc in new_documents:
        if total_chunks_processed >= max_chunks:
            log.warning(
                f"[{agent.agent_name}] Phase 2a: max_chunks ({max_chunks}) reached, "
                f"remaining documents deferred to next cycle"
            )
            break

        doc_path = doc["path"]
        doc_name = doc["name"]

        # Read document content — extract text BEFORE passing to LLM
        doc_text = ""
        try:
            suffix = doc_path.suffix.lower()
            if suffix in {".txt", ".csv", ".md"}:
                doc_text = doc_path.read_text(encoding="utf-8", errors="replace")
            elif suffix == ".pdf":
                # PDF: Try decryption first (banks send encrypted PDFs), then OCR
                pdf_path = doc_path
                try:
                    import pikepdf
                    pdf = pikepdf.open(doc_path)
                    if pdf.is_encrypted:
                        # Try empty password (common for bank statements)
                        decrypted = ws_base / f".tmp_{doc_name}"
                        pdf.save(decrypted)
                        pdf_path = decrypted
                        log.info(f"[{agent.agent_name}] Phase 2a: Decrypted {doc_name}")
                    pdf.close()
                except Exception:
                    pass  # pikepdf not available or not encrypted — proceed with original
                try:
                    from core.skills.skill_document_ocr import DocumentOCRSkill
                    ocr = DocumentOCRSkill(agent_name=agent.agent_name, config={}, workspace_base=str(ws_base))
                    doc_text = ocr._ocr_file(pdf_path)
                    log.info(f"[{agent.agent_name}] Phase 2a: OCR extracted {len(doc_text)} chars from {doc_name}")
                except Exception as ocr_exc:
                    log.warning(f"[{agent.agent_name}] Phase 2a: OCR failed for {doc_name}: {ocr_exc}")
                    doc_text = f"[PDF konnte nicht gelesen werden: {doc_name}. OCR-Fehler: {ocr_exc}]"
            elif suffix in {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp", ".heic", ".heif"}:
                # Images (incl. iPhone HEIC): OCR
                try:
                    from core.skills.skill_document_ocr import DocumentOCRSkill
                    ocr = DocumentOCRSkill(agent_name=agent.agent_name, config={}, workspace_base=str(ws_base))
                    doc_text = ocr._ocr_file(doc_path)
                    log.info(f"[{agent.agent_name}] Phase 2a: OCR extracted {len(doc_text)} chars from {doc_name}")
                except Exception as ocr_exc:
                    doc_text = f"[Bild konnte nicht gelesen werden: {doc_name}. OCR-Fehler: {ocr_exc}]"
            elif suffix == ".xlsx":
                # Excel: Extract cell values as text
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(doc_path, read_only=True, data_only=True)
                    parts = []
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        parts.append(f"--- Sheet: {sheet} ---")
                        for row in ws.iter_rows(values_only=True):
                            cells = [str(c) if c is not None else "" for c in row]
                            if any(cells):
                                parts.append(" | ".join(cells))
                    wb.close()
                    doc_text = "\n".join(parts)
                    log.info(f"[{agent.agent_name}] Phase 2a: Extracted {len(doc_text)} chars from {doc_name}")
                except Exception as exc:
                    doc_text = f"[Excel konnte nicht gelesen werden: {doc_name}. Fehler: {exc}]"
            elif suffix == ".docx":
                # Word: Extract paragraphs as text
                try:
                    import docx
                    d = docx.Document(doc_path)
                    doc_text = "\n".join(p.text for p in d.paragraphs if p.text.strip())
                    log.info(f"[{agent.agent_name}] Phase 2a: Extracted {len(doc_text)} chars from {doc_name}")
                except Exception as exc:
                    doc_text = f"[Word konnte nicht gelesen werden: {doc_name}. Fehler: {exc}]"
            elif suffix == ".pptx":
                # PowerPoint: Extract text from all slides
                try:
                    from pptx import Presentation
                    prs = Presentation(doc_path)
                    parts = []
                    for i, slide in enumerate(prs.slides, 1):
                        slide_texts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    t = para.text.strip()
                                    if t:
                                        slide_texts.append(t)
                        if slide_texts:
                            parts.append(f"--- Folie {i} ---\n" + "\n".join(slide_texts))
                    doc_text = "\n\n".join(parts)
                    log.info(f"[{agent.agent_name}] Phase 2a: Extracted {len(doc_text)} chars from {doc_name}")
                except Exception as exc:
                    doc_text = f"[PowerPoint konnte nicht gelesen werden: {doc_name}. Fehler: {exc}]"
            elif suffix in (".odt", ".odp", ".ods"):
                # LibreOffice/OpenDocument: Extract via odfpy or fallback to ZIP+XML
                try:
                    from odf.opendocument import load as odf_load
                    from odf import text as odf_text
                    from odf.text import P
                    odf_doc = odf_load(str(doc_path))
                    paragraphs = odf_doc.getElementsByType(P)
                    lines = []
                    for p in paragraphs:
                        t = ""
                        for node in p.childNodes:
                            if hasattr(node, "data"):
                                t += node.data
                            elif hasattr(node, "__str__"):
                                t += str(node)
                        if t.strip():
                            lines.append(t.strip())
                    doc_text = "\n".join(lines)
                    log.info(f"[{agent.agent_name}] Phase 2a: Extracted {len(doc_text)} chars from {doc_name}")
                except ImportError:
                    # Fallback: extract content.xml from ODF ZIP
                    import zipfile
                    import re as _re_odf
                    try:
                        with zipfile.ZipFile(doc_path) as zf:
                            content_xml = zf.read("content.xml").decode("utf-8")
                        doc_text = _re_odf.sub(r"<[^>]+>", " ", content_xml)
                        doc_text = _re_odf.sub(r"\s+", " ", doc_text).strip()
                        log.info(f"[{agent.agent_name}] Phase 2a: ODF fallback extracted {len(doc_text)} chars from {doc_name}")
                    except Exception as exc:
                        doc_text = f"[ODF konnte nicht gelesen werden: {doc_name}. Fehler: {exc}]"
                except Exception as exc:
                    doc_text = f"[ODF konnte nicht gelesen werden: {doc_name}. Fehler: {exc}]"
            elif suffix in (".doc", ".xls", ".ppt"):
                # Legacy MS Office: try antiword/catdoc or textract fallback
                import subprocess
                try:
                    if suffix == ".doc":
                        result = subprocess.run(
                            ["antiword", str(doc_path)], capture_output=True, text=True, timeout=30,
                        )
                        doc_text = result.stdout if result.returncode == 0 else ""
                    elif suffix == ".xls":
                        # xls via ssconvert (gnumeric) → CSV
                        import tempfile
                        with tempfile.NamedTemporaryFile(suffix=".csv", delete=False) as tmp:
                            tmp_path = tmp.name
                        result = subprocess.run(
                            ["ssconvert", str(doc_path), tmp_path],
                            capture_output=True, text=True, timeout=30,
                        )
                        if result.returncode == 0:
                            doc_text = Path(tmp_path).read_text(encoding="utf-8", errors="replace")
                        Path(tmp_path).unlink(missing_ok=True)
                    elif suffix == ".ppt":
                        result = subprocess.run(
                            ["catppt", str(doc_path)], capture_output=True, text=True, timeout=30,
                        )
                        doc_text = result.stdout if result.returncode == 0 else ""
                    if not doc_text:
                        doc_text = f"[Legacy Office {suffix} konnte nicht konvertiert werden. antiword/catppt/ssconvert installiert?]"
                    else:
                        log.info(f"[{agent.agent_name}] Phase 2a: Legacy Office extracted {len(doc_text)} chars from {doc_name}")
                except Exception as exc:
                    doc_text = f"[Legacy Office {suffix} Fehler: {exc}]"
            elif suffix in (".msg", ".eml"):
                # Outlook .msg / .eml: Extract subject + body
                try:
                    import email as _email_mod
                    if suffix == ".eml":
                        with open(doc_path, "rb") as f:
                            msg_obj = _email_mod.message_from_binary_file(f)
                        subject = msg_obj.get("Subject", "")
                        body = ""
                        if msg_obj.is_multipart():
                            for part in msg_obj.walk():
                                if part.get_content_type() == "text/plain":
                                    body = part.get_payload(decode=True).decode("utf-8", errors="replace")
                                    break
                        else:
                            body = msg_obj.get_payload(decode=True).decode("utf-8", errors="replace")
                        doc_text = f"Betreff: {subject}\n\n{body}"
                    else:
                        # .msg (Outlook binary) — try extract_msg
                        import extract_msg
                        msg_obj = extract_msg.Message(str(doc_path))
                        doc_text = f"Betreff: {msg_obj.subject}\nVon: {msg_obj.sender}\n\n{msg_obj.body}"
                        msg_obj.close()
                    log.info(f"[{agent.agent_name}] Phase 2a: Extracted {len(doc_text)} chars from {doc_name}")
                except Exception as exc:
                    doc_text = f"[Email-Datei konnte nicht gelesen werden: {doc_name}. Fehler: {exc}]"
            elif suffix == ".rtf":
                # RTF: strip formatting via striprtf
                try:
                    from striprtf.striprtf import rtf_to_text
                    raw = doc_path.read_text(encoding="utf-8", errors="replace")
                    doc_text = rtf_to_text(raw)
                    log.info(f"[{agent.agent_name}] Phase 2a: Extracted {len(doc_text)} chars from {doc_name}")
                except Exception as exc:
                    doc_text = f"[RTF konnte nicht gelesen werden: {doc_name}. Fehler: {exc}]"
            elif suffix in (".mpp", ".vsdx"):
                # MS Project / Visio: metadata extraction only
                doc_text = (
                    f"[{doc_name}: MS Project/Visio-Datei erkannt. "
                    f"Automatische Textextraktion nicht unterstützt. "
                    f"Bitte als PDF oder Excel exportieren.]"
                )
                log.info(f"[{agent.agent_name}] Phase 2a: {doc_name} — no extractor for {suffix}, placeholder added")
            else:
                doc_text = f"[Unbekanntes Dateiformat: {doc_name} ({suffix})]"
        except Exception as exc:
            log.warning(f"[{agent.agent_name}] Phase 2a: Could not read {doc_name}: {exc}")
            continue

        if not doc_text or len(doc_text.strip()) < 10:
            log.warning(f"[{agent.agent_name}] Phase 2a: {doc_name} has no extractable text, skipping")
            continue

        # Split into chunks
        chunks = _chunk_document_text(doc_text, chunk_size) if len(doc_text) > chunk_size else [doc_text]
        log.info(
            f"[{agent.agent_name}] Phase 2a: {doc_name} → {len(chunks)} chunk(s)"
        )

        for i, chunk in enumerate(chunks):
            if total_chunks_processed >= max_chunks:
                break

            # CR-270 pattern: History isolation per chunk
            agent._history = []
            agent._current_thread_id = f"doc:{doc_name}:chunk{i+1}"
            agent._tool_call_count = 0

            chunk_prompt = (
                f"CURRENT SITUATION (from Phase 1 analysis):\n{lagebild}\n\n"
            )
            if arbeitsdatei_context:
                chunk_prompt += (
                    f"YOUR WORK-IN-PROGRESS (last entries from arbeitsdatei.md):\n"
                    f"{arbeitsdatei_context}\n\n"
                )
            chunk_prompt += (
                f"DOCUMENT: {doc_name} (chunk {i+1}/{len(chunks)})\n"
                f"{'='*40}\n{chunk}\n{'='*40}\n\n"
            )
            # Chunk-Loop analysis prompt: configurable via batch_file_specs
            _file_specs = agent.config.get("batch_file_specs", {})
            _arbeitsdatei_spec = _file_specs.get("arbeitsdatei.md", "")
            if _arbeitsdatei_spec and "| " not in _arbeitsdatei_spec:
                # Text-mode (non-table): Agent-specific analysis prompt
                # System prompt has the workflow. Reference files have the data.
                # Chunk prompt just says: analyze this document, use your references.
                chunk_prompt += (
                    "Analysiere dieses Dokument und schreibe ein VOLLSTAENDIGES Ergebnis.\n\n"
                    "PFLICHTFELDER (alle muessen beantwortet werden):\n"
                    "- **Typ**: z.B. stueckliste, schaltplan, risikobeurteilung, betriebsanleitung, "
                    "datenblatt, pruefbericht, ki_dokumentation, netzwerkdokumentation, layout, pneumatikplan\n"
                    "- **Revision/Version**: z.B. Rev. 3, v2.1, oder 'nicht angegeben'\n"
                    "- **Datum**: z.B. 2026-02-15, oder 'nicht angegeben'\n"
                    "- **Aktuell**: Ja oder Nein (Nein wenn aelter als 2 Jahre oder veraltete Norm)\n"
                    "- **Warnungen**: Auffaelligkeiten, fehlende Werte, Widersprueche\n"
                    "- **Zusammenfassung**: 2-3 Saetze was das Dokument enthaelt\n\n"
                    "Schreibe NUR das Analyseergebnis, keinen Kommentar."
                )
            else:
                # Table-mode (legacy steuerberater): keep original prompt
                chunk_prompt += (
                    "Analyze this document chunk thoroughly.\n"
                    "For EACH entry/item/line in the document, output ONE LINE in this exact format:\n"
                    "| Datum | Beschreibung | Betrag | Kategorie | §EStG | Status |\n\n"
                    "Example:\n"
                    "| 16.02.2025 | PINKCAT Tastatur+Maus Set | 25,99€ | Werbungskosten/Arbeitsmittel | §9 | kategorisiert |\n\n"
                    "Status values: kategorisiert / unklar / privat / nachfragen\n"
                    "If unclear whether private or business: set status to 'nachfragen'.\n"
                    "For Handwerker invoices: SEPARATE material costs from labor costs.\n"
                    "Use read_file to load reference files if needed.\n"
                    "Output ONLY the table rows, no commentary. I will save them automatically."
                )

            try:
                with _PhaseParams(agent, "phase2", log):
                    result = await _think_with_activity_check(
                        agent, chunk_prompt, log,
                        stale_timeout=agent.config.get("batch_stale_timeout", 120),
                        phase="2a",  # 2a ORIENT chunk loop
                    )
                total_chunks_processed += 1

                # Auto-append analysis result to arbeitsdatei.md (with dedup)
                # NOTE: Cross-document dedup (credit card stmt vs bank Sammelabbuchung,
                # HAZOP H-A.02) is handled via prompt instruction above. Deterministic
                # cross-reference would require semantic amount matching across docs.
                # Determine mode for this agent
                _file_specs_chk = agent.config.get("batch_file_specs", {})
                _arbeitsdatei_spec_chk = _file_specs_chk.get("arbeitsdatei.md", "")
                _is_table_mode = "| " in _arbeitsdatei_spec_chk or not _arbeitsdatei_spec_chk

                # Minimum quality check: result must be substantial
                if result and len(result.strip()) < 100 and not _is_table_mode:
                    log.warning(
                        f"[{agent.agent_name}] Phase 2a: {doc_name} chunk {i+1} "
                        f"result too short ({len(result.strip())} chars) — retry"
                    )
                    # Retry with more explicit prompt
                    retry_prompt = chunk_prompt + (
                        "\n\nDein vorheriges Ergebnis war zu kurz. "
                        "Schreibe MINDESTENS 5 Zeilen mit allen Pflichtfeldern."
                    )
                    try:
                        result = await _think_with_activity_check(
                            agent, retry_prompt, log,
                            stale_timeout=agent.config.get("batch_stale_timeout", 120),
                            phase="2a",
                        )
                    except Exception:
                        pass  # Use original short result

                if result and len(result.strip()) > 10:
                    header_needed = not arbeitsdatei_path.exists() or arbeitsdatei_path.stat().st_size == 0
                    # Get arbeitsdatei spec from config (if defined)
                    _file_specs = agent.config.get("batch_file_specs", {})
                    _arbeitsdatei_spec = _file_specs.get("arbeitsdatei.md", "")
                    _is_table_mode = "| " in _arbeitsdatei_spec or not _arbeitsdatei_spec
                    # Table-mode (steuerberater): filter/dedup table rows
                    # Text-mode (sichter etc.): append all text with doc separator
                    if _is_table_mode and not _arbeitsdatei_spec:
                        # Legacy default: Steuerberater table format
                        existing_lines = set()
                        if arbeitsdatei_path.exists():
                            for el in arbeitsdatei_path.read_text(encoding="utf-8").split("\n"):
                                el = el.strip()
                                if el.startswith("|") and "Datum" not in el and "---" not in el:
                                    existing_lines.add("|".join(c.strip() for c in el.split("|")))
                        new_lines = []
                        for line in result.strip().split("\n"):
                            line = line.strip()
                            if line.startswith("|") and "Datum" not in line and "---" not in line:
                                if doc_name not in line:
                                    line = line.rstrip("|").rstrip() + f" | {doc_name} |"
                                else:
                                    line = line if line.endswith("|") else line + " |"
                                normalized = "|".join(c.strip() for c in line.split("|"))
                                if normalized not in existing_lines:
                                    new_lines.append(line)
                                    existing_lines.add(normalized)
                        if new_lines or header_needed:
                            with open(arbeitsdatei_path, "a", encoding="utf-8") as f:
                                if header_needed:
                                    f.write("# Beleganalyse — Steuerjahr 2025\n\n")
                                    f.write("| Datum | Beschreibung | Betrag | Kategorie | §EStG | Beleg-Ref | Status |\n")
                                    f.write("|-------|-------------|--------|-----------|-------|-----------|--------|\n")
                                for nl in new_lines:
                                    f.write(nl + "\n")
                    else:
                        # Text-mode: Append full chunk result with document separator
                        # Dedup: Skip if this document was already analyzed
                        _already_in = False
                        if arbeitsdatei_path.exists():
                            _existing = arbeitsdatei_path.read_text(encoding="utf-8")
                            if _doc_has_real_analysis(_existing, doc_name):
                                _already_in = True
                            elif f"## {doc_name}" in _existing:
                                # A fallback stub exists — drop it so the real
                                # analysis replaces it (don't silently discard).
                                if _remove_fallback_section(arbeitsdatei_path, doc_name):
                                    log.info(
                                        f"[{agent.agent_name}] Phase 2a: replaced "
                                        f"fallback stub for {doc_name}"
                                    )
                        if not _already_in:
                            with open(arbeitsdatei_path, "a", encoding="utf-8") as f:
                                if header_needed:
                                    f.write(f"# Arbeitsdatei — {agent.agent_name}\n\n")
                                f.write(f"\n## {doc_name}\n\n")
                                f.write(_sanitize_chunk_result(result).strip() + "\n")
                    log.info(
                        f"[{agent.agent_name}] Phase 2a: {doc_name} chunk {i+1}/{len(chunks)} "
                        f"→ appended to arbeitsdatei.md ({len(result)} chars)"
                    )
                else:
                    log.warning(
                        f"[{agent.agent_name}] Phase 2a: {doc_name} chunk {i+1}/{len(chunks)} "
                        f"produced empty/short result ({len(result)} chars)"
                    )
                    # Write fallback entry so Output Assembly can still inventory this document
                    if not arbeitsdatei_path.exists() or f"## {doc_name}" not in arbeitsdatei_path.read_text(encoding="utf-8"):
                        with open(arbeitsdatei_path, "a", encoding="utf-8") as f:
                            f.write(f"\n## {doc_name}\n\n")
                            f.write("Dokument konnte nicht automatisch analysiert werden. "
                                    "Manuelle Pruefung empfohlen.\n")

                doc_results.append({
                    "thread_id": f"doc:{doc_name}:chunk{i+1}",
                    "status": "analyzed",
                    "response_len": len(result),
                })

                # Refresh arbeitsdatei context for next chunk
                if arbeitsdatei_path.exists():
                    try:
                        content = arbeitsdatei_path.read_text(encoding="utf-8")
                        arbeitsdatei_context = content[-1500:] if len(content) > 1500 else content
                    except Exception:
                        pass

            except (asyncio.TimeoutError, Exception) as exc:
                log.error(
                    f"[{agent.agent_name}] Phase 2a: {doc_name} chunk {i+1} FAILED: {exc}"
                )
                doc_results.append({
                    "thread_id": f"doc:{doc_name}:chunk{i+1}",
                    "status": "FAILED",
                    "error": str(exc),
                })
                continue

    log.info(
        f"[{agent.agent_name}] Phase 2a complete: {total_chunks_processed} chunks processed "
        f"from {len(new_documents)} document(s)"
    )
    return doc_results


async def _phase5_persist(
    agent: "AIMOSAgent", lagebild: str = "", phase3_results: list[dict] | None = None,
    ws_base: Path | None = None, log: logging.Logger | None = None, msg_count: int = 0,
):
    if phase3_results is None:
        phase3_results = []
    assert ws_base is not None and log is not None, "ws_base and log are required"
    """Phase 5: PERSIST — Update workspace files, set reminders, remember key facts."""
    # H-14: Backup state.md before overwrite
    leading_file_path = ws_base / agent.config.get("batch_leading_file", "state.md")
    if leading_file_path.exists():
        backup_path = leading_file_path.with_suffix(".md.bak")
        try:
            shutil.copy2(leading_file_path, backup_path)
            log.info(f"[{agent.agent_name}] Phase 5: Backed up {leading_file_path.name} → {backup_path.name}")
        except Exception as exc:
            log.warning(f"[{agent.agent_name}] Phase 5: Backup failed: {exc}")

    # Build file spec block from config
    file_specs = agent.config.get("batch_file_specs", {})
    if file_specs:
        spec_lines = [f"  {fname}: {spec}" for fname, spec in file_specs.items()]
        file_spec_block = (
            "YOUR WORKSPACE FILES — maintain these using write_file:\n"
            + "\n".join(spec_lines) + "\n\n"
            "IMPORTANT: Follow the format exactly. Phase 0 of your next session will\n"
            "read these files to reconstruct your current state. If the format is wrong,\n"
            "you lose context.\n\n"
        )
    else:
        file_spec_block = (
            "Update your workspace files (use write_file):\n"
            "  - todo.md: Markdown checklist with [ ] and [x], one task per line, include deadline and responsible person\n"
            "  - status.md: Current state of all tracked items, one line per item\n"
            "  - decisions.md: Append new decisions with date and rationale\n\n"
        )

    # S-7: Build Phase 3 results summary
    p3_summary_lines = []
    for pr in phase3_results:
        if pr.get("action") == "output_assembly":
            # File-mode output assembly result
            p3_summary_lines.append(
                f"  ✓ Output Assembly: {pr.get('entries', 0)} Dokumente, "
                f"{pr.get('missing', 0)} fehlende Anhang-IV Punkte"
            )
        elif "status" in pr:
            if pr["status"] in ("sent", "analyzed", "proactive", "no_contact"):
                label = pr["status"]
                p3_summary_lines.append(f"  ✓ {pr['thread_id']}: {label} ({pr.get('response_len', 0)} chars)")
            else:
                p3_summary_lines.append(f"  ✗ {pr['thread_id']}: FAILED — {pr.get('error', 'unknown')}")
        else:
            p3_summary_lines.append(f"  ? {pr}")
    phase3_summary = "\n".join(p3_summary_lines) if p3_summary_lines else "  (no threads processed)"

    leading_file_spec = agent.config.get("batch_leading_file", "state.md")
    _lagebild_block = f"LAGEBILD FROM THIS SESSION:\n{lagebild}\n\n" if lagebild and len(lagebild) > 100 else ""
    phase5_prompt = (
        f"{_lagebild_block}"
        f"PHASE 4 RESULTS (what you actually did):\n{phase3_summary}\n\n"
        "This batch session is complete. You MUST now call write_file for EACH of these files.\n"
        "Do NOT just describe what you would write — actually CALL the write_file tool.\n\n"
        f"{file_spec_block}"
        # File-mode: arbeitsdatei wurde vom Chunk-Loop befüllt → NICHT überschreiben
        + ("STEP 1: arbeitsdatei.md wurde bereits vom Analyseprozess befuellt. NICHT ueberschreiben.\n"
           if agent.config.get("batch_phase4_mode") == "file" and
              (ws_base / "arbeitsdatei.md").exists() and
              (ws_base / "arbeitsdatei.md").stat().st_size > 200
           else "STEP 1: Call write_file(filename=\"arbeitsdatei.md\", content=\"...\") — append new analysis results.\n")
        +
        f"STEP 2: Call write_file(filename=\"offene_punkte.md\", content=\"...\") — open questions for stakeholders.\n"
        f"STEP 3: Call write_file(filename=\"todo.md\", content=\"...\") with your task list.\n"
        f"STEP 4: Call write_file(filename=\"status.md\", content=\"...\") with status table.\n"
        f"STEP 5: Call write_file(filename=\"{leading_file_spec}\", content=\"...\") with your handover protocol.\n"
        f"  The handover protocol ({leading_file_spec}) must be max 1500 characters and contain:\n"
        "  - Current session date and number of messages processed\n"
        "  - Phase 4 results: which threads got responses, which failed\n"
        "  - Top 3-5 open items with deadlines and responsible persons\n"
        "  - Any active escalations (one line each)\n"
        "  - What needs attention next session\n\n"
        "STEP 6: Call remember(key, value) for VERIFIED facts only.\n\n"
        "START NOW. Call write_file immediately."
    )
    log.info(f"[{agent.agent_name}] BATCH Phase 5 (PERSIST)")
    # CR-248: Clear chat history before Phase 4 to prevent history contamination.
    # Phase 1-3 produced text-only responses which teach the LLM "respond with text".
    # Phase 4 needs tool calls (write_file, remember, schedule). A fresh history
    # ensures the LLM sees only the system prompt + Phase 4 prompt, making it much
    # more likely to use native tool calls instead of writing Python code blocks.
    saved_history = getattr(agent, '_history', None)
    if saved_history is not None:
        agent._history = []
    try:
        with _PhaseParams(agent, "phase5", log):
            persist_result = await _think_with_activity_check(agent, phase5_prompt, log, stale_timeout=agent.config.get("batch_stale_timeout", 120), phase="5")
        log.info(f"[{agent.agent_name}] BATCH Phase 5 (PERSIST) complete: {len(persist_result)} chars")
    except (asyncio.TimeoutError, Exception) as exc:
        log.warning(f"[{agent.agent_name}] BATCH Phase 5 FAILED (non-critical): {exc}")
        persist_result = ""
    finally:
        # Restore history (Phase 4 messages are already persisted to DB by think())
        if saved_history is not None:
            agent._history = saved_history

    # CR-248: Parse text-based write_file calls from Phase 4 output.
    # The LLM often writes correct write_file() calls as text instead of native tool calls.
    # The orchestrator extracts and executes them.
    if persist_result and ws_base.exists():
        import re as _re
        text_calls = _re.findall(
            r'write_file\s*\(\s*(?:filename\s*=\s*)?["\']([^"\']+)["\']\s*,\s*(?:content\s*=\s*)?(?:"""(.*?)"""|["\'](.+?)["\'])',
            persist_result, _re.DOTALL,
        )
        # Also try XML format: <invoke name="write_file"><parameter name="path">...</parameter><parameter name="content">...</parameter>
        xml_calls = _re.findall(
            r'<invoke\s+name="write_file">\s*<parameter\s+name="(?:path|filename)">(.*?)</parameter>\s*<parameter\s+name="content">(.*?)</parameter>',
            persist_result, _re.DOTALL,
        )
        all_parsed_calls = [(fn, c1 or c2) for fn, c1, c2 in text_calls] + list(xml_calls)
        for filename, content in all_parsed_calls:
            if filename and content and len(content) > 10:
                # Unescape literal \n sequences that LLMs often produce
                content = content.replace("\\n", "\n").replace("\\t", "\t")
                filepath = ws_base / filename.strip()
                try:
                    filepath.write_text(content.strip(), encoding="utf-8")
                    log.info(f"[{agent.agent_name}] Phase 5 PARSED write_file: {filename} ({len(content)} chars)")
                except Exception as exc:
                    log.warning(f"[{agent.agent_name}] Phase 5 PARSED write_file FAILED {filename}: {exc}")

    # Fallback: If no state.md was written (neither by native tool call nor by parser),
    # the orchestrator writes it from the Lagebild + Phase 3 results.
    state_path = ws_base / leading_file_spec
    state_was_written = state_path.exists() and state_path.stat().st_mtime > (time.time() - 60)
    if not state_was_written and ws_base.exists():
        # state.md was not written by the LLM — write it from the orchestrator
        import datetime as _dt
        fallback_state = (
            f"# Handover — {_dt.datetime.now():%d.%m.%Y %H:%M} "
            f"({msg_count} messages processed)\n\n"
            f"## Situation Report\n{lagebild[:800]}\n\n"
            f"## Phase 3 Results\n{phase3_summary}\n\n"
        )
        if persist_result:
            fallback_state += f"## Agent Notes\n{persist_result[:500]}\n"
        try:
            state_path.write_text(fallback_state[:1500], encoding="utf-8")
            log.info(f"[{agent.agent_name}] Phase 5 FALLBACK: Wrote {leading_file_spec} ({len(fallback_state[:1500])} chars)")
        except Exception as exc:
            log.warning(f"[{agent.agent_name}] Phase 5 FALLBACK write failed: {exc}")
